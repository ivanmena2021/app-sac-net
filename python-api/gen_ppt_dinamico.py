"""
gen_ppt_dinamico.py — Motor de generación de PPT dinámicas para SAC
===================================================================
Genera presentaciones PowerPoint con python-pptx (Python puro).
Incluye: métricas, pipeline SAC, gráficos, tablas, separadores.
Filtros: geográfico, tipo siniestro, empresa, rango de fechas.
"""

import io
import os
import json
import tempfile
import hashlib
import pandas as pd
import numpy as np
from functools import lru_cache
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn


# ══════════════════════════════════════════════════════════════════
# FUNCIONES DE FILTRADO Y CÁLCULO (usadas por app.py para preview)
# ══════════════════════════════════════════════════════════════════

def _safe_col(df, col):
    return col in df.columns


def _safe_num(val, default=0):
    """Sanitize numeric value: convert NaN/inf/None to default."""
    if val is None:
        return default
    try:
        f = float(val)
        if np.isnan(f) or np.isinf(f):
            return default
        return f
    except (ValueError, TypeError):
        return default


def _aplicar_filtros(df, filtros):
    """Aplica todos los filtros al DataFrame (sin copia innecesaria)."""
    mask = pd.Series(True, index=df.index)
    empresa = filtros.get("empresa", "ambas")
    if empresa != "ambas" and _safe_col(df, "EMPRESA"):
        mask &= df["EMPRESA"].astype(str).str.upper().str.contains(empresa.upper(), na=False)
    tipos = filtros.get("tipos_siniestro", [])
    if tipos and _safe_col(df, "TIPO_SINIESTRO"):
        mask &= df["TIPO_SINIESTRO"].isin(tipos)
    fecha_inicio = filtros.get("fecha_inicio")
    fecha_fin = filtros.get("fecha_fin")
    col_fecha = filtros.get("col_fecha", "FECHA_AVISO")
    if fecha_inicio and fecha_fin and _safe_col(df, col_fecha):
        dt_col = pd.to_datetime(df[col_fecha], errors="coerce")
        mask &= (dt_col >= pd.Timestamp(fecha_inicio)) & (dt_col <= pd.Timestamp(fecha_fin))
    deptos = filtros.get("departamentos", [])
    provs = filtros.get("provincias", [])
    dists = filtros.get("distritos", [])
    if deptos and _safe_col(df, "DEPARTAMENTO"):
        mask &= df["DEPARTAMENTO"].isin(deptos)
    if provs and _safe_col(df, "PROVINCIA"):
        mask &= df["PROVINCIA"].isin(provs)
    if dists and _safe_col(df, "DISTRITO"):
        mask &= df["DISTRITO"].isin(dists)
    return df[mask]


def _calcular_metricas(df):
    """Calcula métricas principales del DataFrame (optimizado)."""
    n = len(df)
    if n == 0:
        return {"avisos": 0, "cerrados": 0, "pct_eval": 0, "indemnizacion": 0,
                "desembolso": 0, "pct_desembolso": 0, "ha_indemnizadas": 0, "productores": 0}
    cerrados = int((df["ESTADO_INSPECCION"].astype(str).str.upper() == "CERRADO").sum()) if _safe_col(df, "ESTADO_INSPECCION") else 0
    pct_eval = cerrados / n * 100
    indem = _safe_num(df["INDEMNIZACION"].sum()) if _safe_col(df, "INDEMNIZACION") else 0
    desemb = _safe_num(df["MONTO_DESEMBOLSADO"].sum()) if _safe_col(df, "MONTO_DESEMBOLSADO") else 0
    pct_desemb = (desemb / indem * 100) if indem > 0 else 0
    ha = _safe_num(df["SUP_INDEMNIZADA"].sum()) if _safe_col(df, "SUP_INDEMNIZADA") else 0
    productores = 0
    if _safe_col(df, "N_PRODUCTORES"):
        _prods = pd.to_numeric(df["N_PRODUCTORES"], errors="coerce").fillna(0)
        if _safe_col(df, "INDEMNIZACION"):
            _indemn = pd.to_numeric(df["INDEMNIZACION"], errors="coerce").fillna(0)
            productores = int(_prods[_indemn > 0].sum())
        else:
            productores = int(_prods.sum())
    return {
        "avisos": n, "cerrados": cerrados, "pct_eval": round(pct_eval, 1),
        "indemnizacion": indem, "desembolso": desemb,
        "pct_desembolso": round(pct_desemb, 1),
        "ha_indemnizadas": round(ha, 2), "productores": productores,
    }


def _calcular_pipeline(df):
    """Calcula pipeline del proceso SAC."""
    if not _safe_col(df, "ESTADO_INSPECCION"):
        return []
    raw = df["ESTADO_INSPECCION"].astype(str).str.upper().value_counts().to_dict()
    order = ["NOTIFICADO", "PROGRAMADO", "REPROGRAMADO", "CERRADO"]
    result = []
    for estado in order:
        val = raw.get(estado, 0)
        if val > 0:
            result.append({"label": estado.title(), "val": int(val)})
    for k, v in raw.items():
        if k not in order and v > 0:
            result.append({"label": k.title(), "val": int(v)})
    return result


def _top_breakdown(df, col, n=10):
    """Top N por columna geográfica (optimizado sin copy)."""
    if not _safe_col(df, col):
        return []
    agg = {"Avisos": (col, "count")}
    if _safe_col(df, "INDEMNIZACION"):
        agg["Indemnización"] = ("INDEMNIZACION", "sum")
    if _safe_col(df, "MONTO_DESEMBOLSADO"):
        agg["Desembolso"] = ("MONTO_DESEMBOLSADO", "sum")
    if _safe_col(df, "SUP_INDEMNIZADA"):
        agg["Ha"] = ("SUP_INDEMNIZADA", "sum")

    has_prod = _safe_col(df, "N_PRODUCTORES")
    if has_prod:
        _prods = pd.to_numeric(df["N_PRODUCTORES"], errors="coerce").fillna(0)
        if _safe_col(df, "INDEMNIZACION"):
            _ind = pd.to_numeric(df["INDEMNIZACION"], errors="coerce").fillna(0)
            _prods = _prods.where(_ind > 0, 0)
        prod_by_geo = _prods.groupby(df[col]).sum()

    result = df.groupby(col).agg(**agg).reset_index()
    result = result.sort_values("Avisos", ascending=False).head(n)
    rows = []
    for _, r in result.iterrows():
        row = {"name": str(r[col]), "avisos": int(r["Avisos"])}
        if "Indemnización" in r:
            row["indem"] = float(r["Indemnización"])
        if "Desembolso" in r:
            row["desemb"] = float(r["Desembolso"])
        if "Ha" in r:
            row["ha"] = round(float(r["Ha"]), 2)
        if has_prod and r[col] in prod_by_geo.index:
            row["prod"] = int(prod_by_geo[r[col]])
        rows.append(row)
    return rows


def _tipo_breakdown(df):
    """Distribución por tipo de siniestro."""
    if not _safe_col(df, "TIPO_SINIESTRO"):
        return []
    agg = {"Avisos": ("TIPO_SINIESTRO", "count")}
    if _safe_col(df, "INDEMNIZACION"):
        agg["Indemnización"] = ("INDEMNIZACION", "sum")
    result = df.groupby("TIPO_SINIESTRO").agg(**agg).reset_index()
    result = result.sort_values("Avisos", ascending=False)
    rows = []
    for _, r in result.iterrows():
        row = {"tipo": str(r["TIPO_SINIESTRO"]), "avisos": int(r["Avisos"])}
        if "Indemnización" in r:
            row["indem"] = float(r["Indemnización"])
        rows.append(row)
    return rows


def _empresa_breakdown(df):
    """Breakdown por empresa."""
    if not _safe_col(df, "EMPRESA"):
        return []
    results = []
    for emp in df["EMPRESA"].dropna().unique():
        df_emp = df[df["EMPRESA"] == emp]
        m = _calcular_metricas(df_emp)
        results.append({"empresa": str(emp), **m})
    return results


def _dictamen_breakdown(df):
    """Breakdown por resultado de dictamen."""
    for col in ["DICTAMEN", "RESULTADO_AJUSTE", "RESULTADO_INSPECCION"]:
        if _safe_col(df, col):
            counts = df[col].astype(str).str.upper().value_counts().to_dict()
            return {k: int(v) for k, v in counts.items() if k not in ("NAN", "NONE", "")}
    return {}


LLUVIA_TYPES = {"INUNDACION", "INUNDACIÓN", "HUAYCO", "HUAICO",
                "LLUVIAS EXCESIVAS", "DESLIZAMIENTO", "DESLIZAMIENTOS"}


def _generar_insights(df, metricas, tipos, provincias_o_distritos=None,
                      col_geo="PROVINCIA", provs_seleccionadas=None):
    """Genera insights automáticos basados en los datos."""
    insights = []
    m = metricas
    n = m["avisos"]
    if n == 0:
        return insights

    if tipos and len(tipos) > 0:
        top = tipos[0]
        pct = (top["avisos"] / n * 100) if n > 0 else 0
        indem_txt = f" y {_fmt_money_py(top.get('indem', 0))} en indemnización" if top.get("indem") else ""
        insights.append({
            "title": f"{top['tipo']} es el siniestro predominante",
            "text": f"con {top['avisos']:,} avisos ({pct:.1f}% del total){indem_txt}.",
            "type": "predominance"
        })

    if tipos:
        lluvia = [t for t in tipos if t["tipo"].upper() in LLUVIA_TYPES]
        if lluvia:
            lluvia_avisos = sum(t["avisos"] for t in lluvia)
            lluvia_indem = sum(t.get("indem", 0) for t in lluvia)
            lluvia_nombres = ", ".join(t["tipo"].lower() for t in lluvia[:3])
            insights.append({
                "title": "Eventos asociados a lluvias",
                "text": f"({lluvia_nombres}) suman {lluvia_avisos:,} avisos y {_fmt_money_py(lluvia_indem)}.",
                "type": "lluvia"
            })

    if provincias_o_distritos:
        rezago = [p for p in provincias_o_distritos
                  if p["avisos"] >= 5 and p.get("indem", 0) == 0]
        if rezago:
            top_rez = rezago[0]
            insights.append({
                "title": f"{top_rez['name']} presenta rezago en evaluación",
                "text": f"con {top_rez['avisos']} avisos sin indemnización registrada.",
                "type": "rezago"
            })

    if m["pct_desembolso"] > 0:
        if m["pct_desembolso"] < 30:
            insights.append({
                "title": "Bajo nivel de desembolso",
                "text": f"Solo {m['pct_desembolso']:.1f}% de la indemnización ha sido desembolsada ({_fmt_money_py(m['desembolso'])}).",
                "type": "alert"
            })
        elif m["pct_desembolso"] >= 90:
            insights.append({
                "title": "Alto nivel de desembolso",
                "text": f"{m['pct_desembolso']:.1f}% de la indemnización ya fue desembolsada.",
                "type": "positive"
            })

    if provs_seleccionadas and provincias_o_distritos:
        sel = [p for p in provincias_o_distritos if p["name"] in provs_seleccionadas]
        if sel:
            combined_avisos = sum(p["avisos"] for p in sel)
            combined_indem = sum(p.get("indem", 0) for p in sel)
            combined_ha = sum(p.get("ha", 0) for p in sel)
            combined_prod = sum(p.get("prod", 0) for p in sel)
            nombres = " + ".join(p["name"] for p in sel)
            insights.append({
                "title": f"{nombres}: foco seleccionado",
                "text": f"{combined_avisos:,} avisos combinados, {_fmt_money_py(combined_indem)} indemnización, {combined_ha:,.1f} ha, {combined_prod:,} productores.",
                "type": "highlight"
            })

    return insights[:4]


def _empresa_composition(df):
    """Describe composición de empresa."""
    if not _safe_col(df, "EMPRESA"):
        return ""
    counts = df["EMPRESA"].value_counts()
    total = counts.sum()
    if len(counts) == 1:
        return f"Opera exclusivamente con {counts.index[0]}"
    parts = []
    for emp, cnt in counts.items():
        pct = cnt / total * 100
        parts.append(f"{emp}: {pct:.0f}%")
    return " · ".join(parts)


def _fmt_money_py(n):
    """Format money in Python."""
    if n is None or n == 0:
        return "S/ 0"
    if abs(n) >= 1_000_000:
        return f"S/ {n/1_000_000:,.2f} M"
    return f"S/ {n:,.0f}"


# ══════════════════════════════════════════════════════════════════
# PREPARAR DATA
# ══════════════════════════════════════════════════════════════════

def _prepare_data(df, filtros, fecha_corte):
    """Prepara toda la data en un dict para la generación de slides.

    Modelo de DOS NIVELES ACUMULATIVOS:
      Nivel 1 (Geográfico / obligatorio): resumen total por dimensión geográfica.
        - Se aplican solo filtros de empresa (no tipo_siniestro ni fechas).
        - Nacional → Departamental → Provincial → Distrital según selección.
      Nivel 2 (Complementario / opcional): filtros adicionales (tipos de siniestro,
        fechas de ocurrencia, etc.) que generan secciones complementarias.
        - Si se activa, se aplican los filtros de tipo + fecha sobre el mismo
          alcance geográfico y se agregan slides adicionales.
    """
    scope = filtros.get("scope", "nacional")
    incluir_nacional = filtros.get("incluir_nacional", True)
    deptos = filtros.get("departamentos", [])
    provs = filtros.get("provincias", [])
    dists = filtros.get("distritos", [])

    # ── NIVEL 1: solo filtro de empresa (base geográfica limpia) ──
    filtros_nivel1 = {
        "empresa": filtros.get("empresa", "ambas"),
    }
    df_nivel1 = _aplicar_filtros(df, filtros_nivel1)

    # ── NIVEL 2: empresa + tipos de siniestro + fechas ──
    tipos_sel = filtros.get("tipos_siniestro", [])
    fecha_inicio = filtros.get("fecha_inicio")
    fecha_fin = filtros.get("fecha_fin")
    col_fecha = filtros.get("col_fecha", "FECHA_AVISO")
    hay_nivel2 = bool(tipos_sel) or bool(fecha_inicio and fecha_fin)

    filtros_nivel2 = {
        "empresa": filtros.get("empresa", "ambas"),
        "tipos_siniestro": tipos_sel,
        "fecha_inicio": fecha_inicio,
        "fecha_fin": fecha_fin,
        "col_fecha": col_fecha,
    }
    df_nivel2 = _aplicar_filtros(df, filtros_nivel2) if hay_nivel2 else pd.DataFrame()

    # Etiqueta descriptiva para Nivel 2
    nivel2_label_parts = []
    if tipos_sel:
        nivel2_label_parts.append(f"Tipo: {', '.join(tipos_sel[:3])}")
    if fecha_inicio and fecha_fin:
        nivel2_label_parts.append(f"Período: {fecha_inicio} — {fecha_fin}")
    nivel2_label = " · ".join(nivel2_label_parts) if nivel2_label_parts else ""

    data = {
        "fecha_corte": fecha_corte,
        "scope": scope,
        "filtros": {
            "deptos": deptos, "provs": provs, "dists": dists,
            "tipos": tipos_sel,
            "empresa": filtros.get("empresa", "ambas"),
            "fecha_inicio": str(fecha_inicio) if fecha_inicio else "",
            "fecha_fin": str(fecha_fin) if fecha_fin else "",
        },
        "nivel1_sections": [],
        "nivel2_sections": [],
        "nivel2_label": nivel2_label,
        "hay_nivel2": hay_nivel2,
    }

    # ─────────────────────────────────────────────
    # NIVEL 1: Secciones geográficas (base limpia)
    # ─────────────────────────────────────────────

    if incluir_nacional or not deptos:
        m = _calcular_metricas(df_nivel1)
        tipos_nac = _tipo_breakdown(df_nivel1)
        top_deptos_nac = _top_breakdown(df_nivel1, "DEPARTAMENTO", 10)
        n_deptos = df_nivel1["DEPARTAMENTO"].nunique() if _safe_col(df_nivel1, "DEPARTAMENTO") else 0
        data["nivel1_sections"].append({
            "type": "nacional",
            "metricas": m,
            "pipeline": _calcular_pipeline(df_nivel1),
            "dictamen": _dictamen_breakdown(df_nivel1),
            "empresas": _empresa_breakdown(df_nivel1),
            "top_deptos": top_deptos_nac,
            "tipos": tipos_nac,
            "n_deptos": n_deptos,
            "insights": _generar_insights(df_nivel1, m, tipos_nac, top_deptos_nac, "DEPARTAMENTO"),
        })

    # Pre-agrupar por departamento para evitar filtrados repetidos
    _dept_groups = {}
    if deptos and _safe_col(df_nivel1, "DEPARTAMENTO"):
        for d_name, d_df in df_nivel1[df_nivel1["DEPARTAMENTO"].isin(deptos)].groupby("DEPARTAMENTO"):
            _dept_groups[d_name] = d_df

    if deptos:
        for depto in deptos:
            df_d = _dept_groups.get(depto, pd.DataFrame())
            if len(df_d) == 0:
                continue
            m = _calcular_metricas(df_d)
            tipos_d = _tipo_breakdown(df_d)
            provs_d = _top_breakdown(df_d, "PROVINCIA", 20)
            emp_comp = _empresa_composition(df_d)
            n_provs = df_d["PROVINCIA"].nunique() if _safe_col(df_d, "PROVINCIA") else 0
            data["nivel1_sections"].append({
                "type": "departamental",
                "name": depto,
                "metricas": m,
                "pipeline": _calcular_pipeline(df_d),
                "dictamen": _dictamen_breakdown(df_d),
                "provincias": provs_d,
                "tipos": tipos_d,
                "empresa_comp": emp_comp,
                "n_provincias": n_provs,
                "insights": _generar_insights(df_d, m, tipos_d, provs_d, "PROVINCIA", provs),
                "provs_seleccionadas": provs,
            })

    # Pre-agrupar por provincia
    _prov_groups = {}
    if provs and _safe_col(df_nivel1, "PROVINCIA"):
        for p_name, p_df in df_nivel1[df_nivel1["PROVINCIA"].isin(provs)].groupby("PROVINCIA"):
            _prov_groups[p_name] = p_df

    if provs:
        for prov in provs:
            df_p = _prov_groups.get(prov, pd.DataFrame())
            if len(df_p) == 0:
                continue
            m = _calcular_metricas(df_p)
            depto_name = str(df_p["DEPARTAMENTO"].iloc[0]) if _safe_col(df_p, "DEPARTAMENTO") and len(df_p) > 0 else ""
            tipos_p = _tipo_breakdown(df_p)
            dists_p = _top_breakdown(df_p, "DISTRITO", 20)
            emp_comp = _empresa_composition(df_p)
            data["nivel1_sections"].append({
                "type": "provincial",
                "name": prov,
                "depto": depto_name,
                "metricas": m,
                "pipeline": _calcular_pipeline(df_p),
                "dictamen": _dictamen_breakdown(df_p),
                "distritos": dists_p,
                "tipos": tipos_p,
                "empresa_comp": emp_comp,
                "insights": _generar_insights(df_p, m, tipos_p, dists_p, "DISTRITO"),
            })

    # Pre-agrupar por distrito
    _dist_groups = {}
    if dists and _safe_col(df_nivel1, "DISTRITO"):
        for dt_name, dt_df in df_nivel1[df_nivel1["DISTRITO"].isin(dists[:5])].groupby("DISTRITO"):
            _dist_groups[dt_name] = dt_df

    if dists:
        for dist in dists[:5]:
            df_dist = _dist_groups.get(dist, pd.DataFrame())
            if len(df_dist) == 0:
                continue
            m = _calcular_metricas(df_dist)
            prov_name = str(df_dist["PROVINCIA"].iloc[0]) if _safe_col(df_dist, "PROVINCIA") and len(df_dist) > 0 else ""
            depto_name = str(df_dist["DEPARTAMENTO"].iloc[0]) if _safe_col(df_dist, "DEPARTAMENTO") and len(df_dist) > 0 else ""
            data["nivel1_sections"].append({
                "type": "distrital",
                "name": dist,
                "prov": prov_name,
                "depto": depto_name,
                "metricas": m,
                "pipeline": _calcular_pipeline(df_dist),
                "tipos": _tipo_breakdown(df_dist),
            })

    # ─────────────────────────────────────────────
    # NIVEL 2: Secciones complementarias (filtradas)
    # ─────────────────────────────────────────────

    if hay_nivel2 and len(df_nivel2) > 0:
        # Nacional filtrado
        if incluir_nacional or not deptos:
            m2 = _calcular_metricas(df_nivel2)
            tipos_n2 = _tipo_breakdown(df_nivel2)
            top_deptos_n2 = _top_breakdown(df_nivel2, "DEPARTAMENTO", 10)
            data["nivel2_sections"].append({
                "type": "nacional",
                "metricas": m2,
                "pipeline": _calcular_pipeline(df_nivel2),
                "tipos": tipos_n2,
                "top_deptos": top_deptos_n2,
                "empresas": _empresa_breakdown(df_nivel2),
                "insights": _generar_insights(df_nivel2, m2, tipos_n2, top_deptos_n2, "DEPARTAMENTO"),
            })

        # Departamentos filtrados (pre-agrupado)
        _dept_groups2 = {}
        if deptos and _safe_col(df_nivel2, "DEPARTAMENTO"):
            for d2n, d2f in df_nivel2[df_nivel2["DEPARTAMENTO"].isin(deptos)].groupby("DEPARTAMENTO"):
                _dept_groups2[d2n] = d2f

        if deptos:
            for depto in deptos:
                df_d2 = _dept_groups2.get(depto, pd.DataFrame())
                if len(df_d2) == 0:
                    continue
                m2 = _calcular_metricas(df_d2)
                tipos_d2 = _tipo_breakdown(df_d2)
                provs_d2 = _top_breakdown(df_d2, "PROVINCIA", 20)
                data["nivel2_sections"].append({
                    "type": "departamental",
                    "name": depto,
                    "metricas": m2,
                    "pipeline": _calcular_pipeline(df_d2),
                    "provincias": provs_d2,
                    "tipos": tipos_d2,
                    "empresa_comp": _empresa_composition(df_d2),
                    "insights": _generar_insights(df_d2, m2, tipos_d2, provs_d2, "PROVINCIA", provs),
                })

        # Provincias filtradas (pre-agrupado)
        _prov_groups2 = {}
        if provs and _safe_col(df_nivel2, "PROVINCIA"):
            for p2n, p2f in df_nivel2[df_nivel2["PROVINCIA"].isin(provs)].groupby("PROVINCIA"):
                _prov_groups2[p2n] = p2f

        if provs:
            for prov in provs:
                df_p2 = _prov_groups2.get(prov, pd.DataFrame())
                if len(df_p2) == 0:
                    continue
                m2 = _calcular_metricas(df_p2)
                depto_name = str(df_p2["DEPARTAMENTO"].iloc[0]) if _safe_col(df_p2, "DEPARTAMENTO") and len(df_p2) > 0 else ""
                tipos_p2 = _tipo_breakdown(df_p2)
                dists_p2 = _top_breakdown(df_p2, "DISTRITO", 20)
                data["nivel2_sections"].append({
                    "type": "provincial",
                    "name": prov,
                    "depto": depto_name,
                    "metricas": m2,
                    "distritos": dists_p2,
                    "tipos": tipos_p2,
                    "insights": _generar_insights(df_p2, m2, tipos_p2, dists_p2, "DISTRITO"),
                })

    return data


# ══════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS FOR PPT GENERATION (python-pptx)
# ══════════════════════════════════════════════════════════════════

# ── Paleta institucional MIDAGRI ──
C = {
    "forest": RGBColor(0x40, 0x8B, 0x14),      # Verde MIDAGRI principal (#408B14)
    "green": RGBColor(0x52, 0xB0, 0x17),        # Verde claro indicadores (#52B017)
    "sage": RGBColor(0x52, 0xB0, 0x17),          # Verde claro (alias)
    "mint": RGBColor(0x95, 0xD5, 0xB2),          # Verde suave
    "cream": RGBColor(0xF2, 0xF2, 0xF2),        # Gris claro fondo MIDAGRI (#F2F2F2)
    "gold": RGBColor(0xFF, 0xC0, 0x00),          # Dorado MIDAGRI (#FFC000)
    "amber": RGBColor(0xFF, 0xC0, 0x00),         # Dorado (alias)
    "navy": RGBColor(0x3F, 0x3F, 0x3F),          # Gris oscuro MIDAGRI (#3F3F3F)
    "dark": RGBColor(0x21, 0x25, 0x29),          # Negro profundo
    "gray": RGBColor(0x59, 0x59, 0x59),          # Gris medio MIDAGRI (#595959)
    "lightGray": RGBColor(0xD8, 0xD8, 0xD8),    # Gris claro MIDAGRI (#D8D8D8)
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "red": RGBColor(0xC0, 0x39, 0x2B),
    "blue": RGBColor(0x21, 0x96, 0xF3),
    "orange": RGBColor(0xFF, 0xC0, 0x00),        # Dorado MIDAGRI (reemplaza naranja)
    "coral": RGBColor(0xE7, 0x4C, 0x3C),
    "teal": RGBColor(0x10, 0xA9, 0xA7),          # Teal MIDAGRI (#10A9A7)
    "yellowBg": RGBColor(0xFF, 0xF3, 0xCD),
    "lightCream": RGBColor(0xF8, 0xF9, 0xFA),
}


# ── Corrección ortográfica de nombres geográficos peruanos ──
_GEO_TILDES = {
    # Departamentos
    "AMAZONAS": "Amazonas", "ANCASH": "Áncash", "APURIMAC": "Apurímac",
    "AREQUIPA": "Arequipa", "AYACUCHO": "Ayacucho", "CAJAMARCA": "Cajamarca",
    "CALLAO": "Callao", "CUSCO": "Cusco", "HUANCAVELICA": "Huancavelica",
    "HUANUCO": "Huánuco", "ICA": "Ica", "JUNIN": "Junín",
    "LA LIBERTAD": "La Libertad", "LAMBAYEQUE": "Lambayeque", "LIMA": "Lima",
    "LORETO": "Loreto", "MADRE DE DIOS": "Madre de Dios",
    "MOQUEGUA": "Moquegua", "PASCO": "Pasco", "PIURA": "Piura",
    "PUNO": "Puno", "SAN MARTIN": "San Martín", "TACNA": "Tacna",
    "TUMBES": "Tumbes", "UCAYALI": "Ucayali",
    # Provincias con tilde frecuente
    "MOYOBAMBA": "Moyobamba", "RIOJA": "Rioja", "LAMAS": "Lamas",
    "SAN JOSE DE SISA": "San José de Sisa", "MARISCAL CACERES": "Mariscal Cáceres",
    "HUALLAGA": "Huallaga", "BELLAVISTA": "Bellavista", "PICOTA": "Picota",
    "TOCACHE": "Tocache", "EL DORADO": "El Dorado",
    "CONCEPCION": "Concepción", "JAEN": "Jaén", "SATIPO": "Satipo",
    "HUAMANGA": "Huamanga", "CAÑETE": "Cañete", "BARRANCA": "Barranca",
    "HUAURA": "Huaura", "MAYNAS": "Maynas", "CORONEL PORTILLO": "Coronel Portillo",
    "SANCHEZ CARRION": "Sánchez Carrión", "SANTIAGO DE CHUCO": "Santiago de Chuco",
    "BOLIVAR": "Bolívar", "PARINACOCHAS": "Parinacochas",
    "VICTOR FAJARDO": "Víctor Fajardo", "PABON": "Pabón",
    "RODRIGUEZ DE MENDOZA": "Rodríguez de Mendoza",
    "GRAN CHIMU": "Gran Chimú", "CUTERVO": "Cutervo",
    "CHOTA": "Chota", "SANTA CRUZ": "Santa Cruz",
    "SAN IGNACIO": "San Ignacio", "SAN MARCOS": "San Marcos",
    "SAN MIGUEL": "San Miguel", "SAN PABLO": "San Pablo",
    "CONTUMAZA": "Contumazá", "CELENDIN": "Celendín",
    "HUALGAYOC": "Hualgayoc", "CHICLAYO": "Chiclayo",
    "FERREÑAFE": "Ferreñafe", "TRUJILLO": "Trujillo",
    "ASCOPE": "Ascope", "CHEPEN": "Chepén",
    "PACASMAYO": "Pacasmayo", "PATAZ": "Pataz",
    "OTUZCO": "Otuzco", "VIRU": "Virú",
}


@lru_cache(maxsize=512)
def _fix_geo_name(name):
    """Corregir tildes en nombres geográficos peruanos."""
    if not name:
        return name
    key = name.strip().upper()
    if key in _GEO_TILDES:
        return _GEO_TILDES[key]
    # Si no está en diccionario, usar title() como fallback
    return name.strip().title()


@lru_cache(maxsize=512)
def _fix_geo_upper(name):
    """Versión en mayúsculas con tildes correctas."""
    return _fix_geo_name(name).upper()


def _fmt_num(n):
    """Format number with locale."""
    if n is None:
        return "0"
    return f"{int(n):,}"


def _fmt_money(n):
    """Format money (S/ X,XXX or S/ X.XX M)."""
    if n is None or n == 0:
        return "S/ 0"
    if abs(n) >= 1_000_000:
        return f"S/{n/1_000_000:,.2f}M"
    if abs(n) >= 100_000:
        return f"S/{n/1_000:,.0f}K"
    return f"S/ {n:,.0f}"


def _fmt_pct(n):
    """Format percentage."""
    if n is None:
        return "0%"
    return f"{n:.1f}%"


def _add_shadow(shape, alpha_val=15000):
    """Add shadow effect to shape via XML."""
    try:
        spPr = shape._element.spPr
        effectLst = spPr.makeelement(qn('a:effectLst'), {})
        outerShdw = effectLst.makeelement(qn('a:outerShdw'), {
            'blurRad': '50800',
            'dist': '25400',
            'dir': '8100000',
        })
        srgbClr = outerShdw.makeelement(qn('a:srgbClr'), {'val': '000000'})
        alpha = srgbClr.makeelement(qn('a:alpha'), {'val': str(alpha_val)})
        srgbClr.append(alpha)
        outerShdw.append(srgbClr)
        effectLst.append(outerShdw)
        spPr.append(effectLst)
    except Exception:
        pass


def _make_logo(slide, x, y, w, h):
    """Create MIDAGRI/SAC logo using shapes: white rounded rect with two teal rects inside."""
    logo_bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x, y, w, h
    )
    logo_bg.fill.solid()
    logo_bg.fill.fore_color.rgb = C["white"]
    logo_bg.line.color.rgb = C["teal"]
    logo_bg.line.width = Pt(2)

    rect_w = w * 0.35
    rect_h = h * 0.6
    gap = (w - rect_w * 2) / 3

    rect1 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x + gap, y + (h - rect_h) / 2,
        rect_w, rect_h
    )
    rect1.fill.solid()
    rect1.fill.fore_color.rgb = C["teal"]
    rect1.line.fill.background()

    rect2 = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        x + gap * 2 + rect_w, y + (h - rect_h) / 2,
        rect_w, rect_h
    )
    rect2.fill.solid()
    rect2.fill.fore_color.rgb = C["sage"]
    rect2.line.fill.background()


def _add_kpi_card(slide, left, top, w, h, label, value, sublabel, accent_color, icon_text):
    """Add enhanced metric card with colored top bar, icon circle, and large number."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, w, h
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = C["white"]
    bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    bg.line.width = Pt(1)
    _add_shadow(bg, 15000)

    top_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, w, Inches(0.25)
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent_color
    top_bar.line.fill.background()

    icon_circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        left + (w - Inches(0.45)) / 2, top + Inches(0.35),
        Inches(0.45), Inches(0.45)
    )
    icon_circle.fill.solid()
    icon_circle.fill.fore_color.rgb = accent_color
    icon_circle.line.fill.background()

    icon_tf = slide.shapes.add_textbox(
        left + (w - Inches(0.45)) / 2, top + Inches(0.35),
        Inches(0.45), Inches(0.45)
    )
    icon_frame = icon_tf.text_frame
    icon_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    icon_p = icon_frame.paragraphs[0]
    icon_p.text = icon_text
    icon_p.font.size = Pt(20)
    icon_p.alignment = PP_ALIGN.CENTER

    # Auto-size value font based on text length
    val_str = str(value)
    if len(val_str) > 10:
        val_font = Pt(18)
    elif len(val_str) > 7:
        val_font = Pt(22)
    else:
        val_font = Pt(26)

    value_tf = slide.shapes.add_textbox(
        left + Inches(0.08), top + Inches(0.88),
        w - Inches(0.16), Inches(0.45)
    )
    value_frame = value_tf.text_frame
    value_frame.word_wrap = False
    value_frame.margin_top = Pt(0)
    value_frame.margin_bottom = Pt(0)
    value_p = value_frame.paragraphs[0]
    value_p.text = val_str
    value_p.font.size = val_font
    value_p.font.bold = True
    value_p.font.color.rgb = C["navy"]
    value_p.font.name = "Calibri"
    value_p.alignment = PP_ALIGN.CENTER

    sub_tf = slide.shapes.add_textbox(
        left + Inches(0.08), top + Inches(1.35),
        w - Inches(0.16), Inches(0.25)
    )
    sub_frame = sub_tf.text_frame
    sub_frame.word_wrap = True
    sub_frame.margin_top = Pt(0)
    sub_frame.margin_bottom = Pt(0)
    sub_p = sub_frame.paragraphs[0]
    sub_p.text = str(sublabel) if sublabel else ""
    sub_p.font.size = Pt(8)
    sub_p.font.color.rgb = C["gray"]
    sub_p.alignment = PP_ALIGN.CENTER

    label_tf = slide.shapes.add_textbox(
        left + Inches(0.08), top + Inches(1.58),
        w - Inches(0.16), Inches(0.22)
    )
    label_frame = label_tf.text_frame
    label_frame.word_wrap = True
    label_frame.margin_top = Pt(0)
    label_frame.margin_bottom = Pt(0)
    label_p = label_frame.paragraphs[0]
    label_p.text = str(label)
    label_p.font.size = Pt(9)
    label_p.font.color.rgb = C["gray"]
    label_p.alignment = PP_ALIGN.CENTER


def _add_header_bar(slide, title, color, y_pos=Inches(0.3)):
    """Add colored rectangle header bar with title."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), y_pos,
        Inches(10), Inches(0.7)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = color

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C["white"]
    p.alignment = PP_ALIGN.LEFT
    text_frame.margin_left = Inches(0.3)
    text_frame.margin_top = Inches(0.1)
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_styled_table(slide, headers, rows, left=Inches(0.3), top=Inches(1.2),
                     col_widths=None, has_total=False, max_rows=12):
    """Add professional table with alternating rows, header styling, optional total row."""
    rows_to_add = min(len(rows), max_rows)
    if has_total:
        rows_to_add += 1
    cols = len(headers)

    table_shape = slide.shapes.add_table(rows_to_add + 1, cols, left, top,
                                         Inches(9.4), Inches(0.35 * (rows_to_add + 1)))
    table = table_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C["navy"]
        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = C["white"]
        p.alignment = PP_ALIGN.CENTER

    for row_idx, row_data in enumerate(rows[:max_rows], 1):
        for col_idx, value in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.fill.solid()
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = C["white"]
            else:
                cell.fill.fore_color.rgb = C["lightCream"]

            p = cell.text_frame.paragraphs[0]
            p.text = str(value) if value is not None else ""
            p.font.size = Pt(9)
            p.font.color.rgb = C["dark"]
            p.alignment = PP_ALIGN.CENTER

    if has_total and len(rows) > 0:
        total_row = rows_to_add
        for col_idx in range(cols):
            cell = table.cell(total_row, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C["navy"]
            p = cell.text_frame.paragraphs[0]
            p.text = "TOTAL" if col_idx == 0 else ""
            p.font.size = Pt(9)
            p.font.bold = True
            p.font.color.rgb = C["white"]
            p.alignment = PP_ALIGN.CENTER


def _add_pipeline_flow(slide, pipeline, y_pos):
    """Add N-step pipeline flow — adapta tamaño de tarjetas al número de pasos."""
    if not pipeline:
        return

    n_steps = min(len(pipeline), 5)
    total = sum(p["val"] for p in pipeline)
    colors = [C["teal"], C["amber"], C["teal"], C["navy"], C["orange"]]

    # Calcular dimensiones dinámicas según cantidad de pasos
    available_w = 9.2  # pulgadas útiles (0.4 a 9.6)
    gap = 0.15  # espacio entre tarjetas (para flechas)
    card_w_val = (available_w - gap * (n_steps - 1)) / n_steps
    card_w = Inches(card_w_val)
    card_h = Inches(2.6)
    step_w = card_w_val + gap
    circle_size = Inches(0.8)
    x_start = Inches(0.4)

    for i, stage in enumerate(pipeline[:n_steps]):
        x_pos = x_start + Inches(i * step_w)

        card_bg = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x_pos, y_pos, card_w, card_h
        )
        card_bg.fill.solid()
        card_bg.fill.fore_color.rgb = C["white"]
        card_bg.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        card_bg.line.width = Pt(1)
        _add_shadow(card_bg, 12000)

        top_bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x_pos, y_pos, card_w, Inches(0.2)
        )
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = colors[i % len(colors)]
        top_bar.line.fill.background()

        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            x_pos + (card_w - circle_size) / 2, y_pos + Inches(0.35),
            circle_size, circle_size
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = colors[i % len(colors)]
        circle.line.fill.background()

        circle_tf = slide.shapes.add_textbox(
            x_pos + (card_w - circle_size) / 2, y_pos + Inches(0.35),
            circle_size, circle_size
        )
        circle_frame = circle_tf.text_frame
        circle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        circle_p = circle_frame.paragraphs[0]
        circle_p.text = str(i + 1)
        circle_p.font.size = Pt(26)
        circle_p.font.bold = True
        circle_p.font.color.rgb = C["white"]
        circle_p.alignment = PP_ALIGN.CENTER

        label_tf = slide.shapes.add_textbox(
            x_pos + Inches(0.05), y_pos + Inches(1.3),
            card_w - Inches(0.1), Inches(0.35)
        )
        label_frame = label_tf.text_frame
        label_frame.word_wrap = True
        label_p = label_frame.paragraphs[0]
        label_p.text = stage["label"]
        label_p.font.size = Pt(10)
        label_p.font.color.rgb = C["gray"]
        label_p.alignment = PP_ALIGN.CENTER

        val_tf = slide.shapes.add_textbox(
            x_pos + Inches(0.05), y_pos + Inches(1.7),
            card_w - Inches(0.1), Inches(0.35)
        )
        val_frame = val_tf.text_frame
        val_frame.word_wrap = False
        val_p = val_frame.paragraphs[0]
        val_p.text = f"{stage['val']:,}"
        val_p.font.size = Pt(20)
        val_p.font.bold = True
        val_p.font.color.rgb = C["navy"]
        val_p.alignment = PP_ALIGN.CENTER

        pct = (stage["val"] / total * 100) if total > 0 else 0
        pct_tf = slide.shapes.add_textbox(
            x_pos + Inches(0.05), y_pos + Inches(2.1),
            card_w - Inches(0.1), Inches(0.25)
        )
        pct_frame = pct_tf.text_frame
        pct_frame.word_wrap = True
        pct_p = pct_frame.paragraphs[0]
        pct_p.text = f"{pct:.1f}%"
        pct_p.font.size = Pt(12)
        pct_p.font.bold = True
        pct_p.font.color.rgb = colors[i % len(colors)]
        pct_p.alignment = PP_ALIGN.CENTER

        # Flecha entre tarjetas
        if i < n_steps - 1:
            arrow_x = x_pos + card_w + Inches(0.01)
            arrow_y = y_pos + Inches(1.2)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW,
                arrow_x, arrow_y,
                Inches(gap - 0.02), Inches(0.22)
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = C["lightGray"]
            arrow.line.fill.background()


def _add_alert_box(slide, text, left, top, width):
    """Add yellow/amber alert box with warning content."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, Inches(0.55)
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = C["yellowBg"]
    bg.line.color.rgb = C["amber"]
    bg.line.width = Pt(1.5)

    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, Inches(0.05), Inches(0.55)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = C["amber"]
    accent_bar.line.fill.background()

    tf = slide.shapes.add_textbox(
        left + Inches(0.15), top + Inches(0.08),
        width - Inches(0.3), Inches(0.4)
    )
    text_frame = tf.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "⚠ " + text
    p.font.size = Pt(10)
    p.font.color.rgb = C["dark"]
    p.alignment = PP_ALIGN.LEFT


def _add_resumen_ejecutivo_slide(slide, prs, dept_name, text, fecha_corte):
    """Add dark background resumen ejecutivo slide with narrative text."""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["navy"]
    background.line.fill.background()

    line_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.2),
        prs.slide_width, Inches(0.06)
    )
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = C["teal"]
    line_top.line.fill.background()

    line_bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.42),
        prs.slide_width, Inches(0.06)
    )
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = C["teal"]
    line_bottom.line.fill.background()

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(0.4),
        Inches(2.0), Inches(0.35)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = C["white"]
    badge.line.color.rgb = C["teal"]
    badge.line.width = Pt(1)

    badge_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.4),
        Inches(2.0), Inches(0.35)
    )
    badge_frame = badge_tf.text_frame
    badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    badge_p = badge_frame.paragraphs[0]
    badge_p.text = _fix_geo_upper(dept_name)
    badge_p.font.size = Pt(12)
    badge_p.font.bold = True
    badge_p.font.color.rgb = C["teal"]
    badge_p.alignment = PP_ALIGN.CENTER

    title_tf = slide.shapes.add_textbox(
        Inches(2.6), Inches(0.35),
        Inches(6.8), Inches(0.45)
    )
    title_frame = title_tf.text_frame
    title_frame.word_wrap = True
    title_p = title_frame.paragraphs[0]
    title_p.text = "Resumen Ejecutivo"
    title_p.font.size = Pt(24)
    title_p.font.bold = True
    title_p.font.color.rgb = C["lightGray"]
    title_p.font.name = "Calibri"
    title_p.alignment = PP_ALIGN.LEFT

    sep_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.4), Inches(0.85),
        Inches(9.2), Inches(0.02)
    )
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = C["teal"]
    sep_line.line.fill.background()

    body_tf = slide.shapes.add_textbox(
        Inches(0.6), Inches(1.1),
        Inches(8.8), Inches(4.15)
    )
    body_frame = body_tf.text_frame
    body_frame.word_wrap = True
    body_p = body_frame.paragraphs[0]
    body_p.text = text
    body_p.font.size = Pt(15)
    body_p.font.color.rgb = C["white"]
    body_p.font.name = "Calibri"
    body_p.alignment = PP_ALIGN.JUSTIFY
    body_p.line_spacing = 1.4

    footer_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(5.2),
        Inches(9.2), Inches(0.2)
    )
    footer_frame = footer_tf.text_frame
    footer_frame.word_wrap = True
    footer_p = footer_frame.paragraphs[0]
    footer_p.text = f"Fuente: DSFFA — MIDAGRI, SAC 2025-2026. Corte al {fecha_corte}."
    footer_p.font.size = Pt(9)
    footer_p.font.italic = True
    footer_p.font.color.rgb = C["lightGray"]
    footer_p.alignment = PP_ALIGN.RIGHT


# ══════════════════════════════════════════════════════════════════
# SLIDE GENERATION FUNCTIONS
# ══════════════════════════════════════════════════════════════════

def _generar_resumen_texto(section_data, scope_type):
    """Generate narrative paragraph text from section metrics — texto completo."""
    m = section_data.get("metricas", {})
    tipos = section_data.get("tipos", [])
    provs = section_data.get("provincias", [])

    nombre = section_data.get("name", "La región")
    avisos = m.get("avisos", 0)
    pct_eval = m.get("pct_eval", 0)
    cerrados = m.get("cerrados", 0)
    indem_total = m.get("indemnizacion", 0)
    pct_desembolso = m.get("pct_desembolso", 0)
    desembolso = m.get("desembolso", 0)
    ha = m.get("ha_indemnizadas", 0)
    productores = m.get("productores", 0)
    pendientes = avisos - cerrados

    tipo_predominante = tipos[0]["tipo"] if tipos else "eventos registrados"
    tipo_pct = (tipos[0]["avisos"] / avisos * 100) if tipos and avisos > 0 else 0
    prov_principal = provs[0]["name"] if provs else "las zonas con mayor concentración"

    # Segundo tipo si existe
    segundo_tipo = ""
    if len(tipos) > 1:
        segundo_tipo = f", seguido de {tipos[1]['tipo']} con {tipos[1]['avisos']:,} avisos"

    # Construir párrafos temáticos
    parrafo_general = (
        f"{nombre} registra {avisos:,} avisos de siniestro con un avance de evaluación "
        f"del {pct_eval:.1f}% ({cerrados:,} expedientes cerrados de {avisos:,} reportados). "
    )

    parrafo_financiero = (
        f"La indemnización total reconocida asciende a S/ {indem_total:,.0f}, con un "
        f"desembolso acumulado de S/ {desembolso:,.0f} ({pct_desembolso:.1f}% del monto "
        f"reconocido). Se han beneficiado {productores:,} productores sobre "
        f"{ha:,.1f} hectáreas indemnizadas. "
    )

    parrafo_tipos = (
        f"El tipo de siniestro predominante es {tipo_predominante}, que representa "
        f"el {tipo_pct:.1f}% del total de avisos{segundo_tipo}. La mayor carga de "
        f"siniestros se concentra en {prov_principal}. "
    )

    parrafo_pendientes = ""
    if pendientes > 0:
        parrafo_pendientes = (
            f"Quedan {pendientes:,} expedientes pendientes de cierre, lo que requiere "
            f"acelerar los procesos de inspección y dictamen para cumplir con los "
            f"compromisos de desembolso hacia los productores afectados."
        )

    text = parrafo_general + parrafo_financiero + parrafo_tipos + parrafo_pendientes
    return text


def _add_pipeline_slide(prs, pipeline, dictamen, metricas):
    """Add pipeline process flow slide — sin alerta, centrado verticalmente."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["cream"]
    background.line.fill.background()

    tf_title = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.35),
        Inches(9.2), Inches(0.45)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Flujo de Procesos del SAC — Estado Actual"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C["navy"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.85),
        prs.slide_width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C["teal"]
    line.line.fill.background()

    # Centrar pipeline verticalmente en el espacio disponible (0.95 a 5.4)
    _add_pipeline_flow(slide, pipeline, Inches(1.5))


def _add_resumen_ejecutivo(prs, section_name, resumen_text, scope_label, fecha_corte):
    """Add resumen ejecutivo slide with dark background and narrative text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_resumen_ejecutivo_slide(slide, prs, section_name, resumen_text, fecha_corte)


def _add_tipo_siniestro_slide(prs, tipos):
    """Add tipo de siniestro: tabla compacta a la izquierda + chart de barras a la derecha."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["cream"]
    background.line.fill.background()

    tf_title = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.35),
        Inches(9.2), Inches(0.45)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Distribución por Tipo de Siniestro"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C["navy"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.85),
        prs.slide_width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C["teal"]
    line.line.fill.background()

    # ── Tabla compacta izquierda (con Avisos) ──
    headers = ["Tipo Siniestro", "Avisos", "Indem. (S/)", "%"]
    rows = []
    total_monto = sum(t.get("indem", 0) for t in tipos)
    total_avisos = sum(t.get("avisos", 0) for t in tipos)

    for t in tipos[:10]:
        monto = t.get("indem", 0)
        pct = (monto / total_monto * 100) if total_monto > 0 else 0
        rows.append([
            t["tipo"],
            _fmt_num(t.get("avisos", 0)),
            _fmt_money(monto),
            f"{pct:.1f}%"
        ])

    col_widths = [Inches(2.2), Inches(0.8), Inches(1.2), Inches(0.7)]
    _add_styled_table(slide, headers, rows, left=Inches(0.3), top=Inches(1.15),
                     col_widths=col_widths, max_rows=10)

    # ── Chart de barras horizontales a la derecha ──
    top_tipos = tipos[:8]
    if top_tipos:
        chart_data = CategoryChartData()
        chart_data.categories = [t["tipo"][:18] for t in top_tipos]
        chart_data.add_series('Indemnización', tuple(_safe_num(t.get("indem", 0)) for t in top_tipos))

        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.BAR_CLUSTERED, Inches(5.3), Inches(1.15),
            Inches(4.4), Inches(3.8), chart_data
        ).chart

        chart_shape.has_legend = False
        chart_shape.series[0].format.fill.solid()
        chart_shape.series[0].format.fill.fore_color.rgb = C["teal"]

        # Etiquetas de datos
        try:
            chart_shape.series[0].has_data_labels = True
            data_labels = chart_shape.series[0].data_labels
            data_labels.font.size = Pt(8)
            data_labels.font.color.rgb = C["navy"]
        except Exception:
            pass

    # ── Resumen rápido ──
    summary_tf = slide.shapes.add_textbox(
        Inches(5.3), Inches(5.05),
        Inches(4.4), Inches(0.25)
    )
    summary_frame = summary_tf.text_frame
    summary_frame.word_wrap = True
    summary_p = summary_frame.paragraphs[0]
    summary_p.text = f"Total: {total_avisos:,} avisos · {_fmt_money(total_monto)} indemnización"
    summary_p.font.size = Pt(9)
    summary_p.font.bold = True
    summary_p.font.color.rgb = C["navy"]
    summary_p.alignment = PP_ALIGN.CENTER

    footer_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(5.35),
        Inches(9.2), Inches(0.2)
    )
    footer_frame = footer_tf.text_frame
    footer_frame.word_wrap = True
    footer_p = footer_frame.paragraphs[0]
    footer_p.text = "Fuente: DSFFA — MIDAGRI, SAC 2025-2026"
    footer_p.font.size = Pt(8)
    footer_p.font.italic = True
    footer_p.font.color.rgb = C["gray"]
    footer_p.alignment = PP_ALIGN.RIGHT


def _add_top_deptos_chart(prs, top_deptos):
    """Add top departamentos bar chart slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["cream"]
    background.line.fill.background()

    tf_title = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.35),
        Inches(9.2), Inches(0.45)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Top 12 Departamentos por Indemnización Reconocida"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = C["navy"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.85),
        prs.slide_width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C["teal"]
    line.line.fill.background()

    # Create chart data
    chart_data = CategoryChartData()
    chart_data.categories = [d["name"] for d in top_deptos[:12]]

    values = [_safe_num(d.get("indem", 0)) for d in top_deptos[:12]]
    chart_data.add_series('Indemnización (S/)', tuple(values))

    # Add chart
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.3),
        Inches(9.0), Inches(3.8), chart_data
    ).chart

    chart_shape.has_legend = False
    chart_shape.series[0].format.fill.solid()
    chart_shape.series[0].format.fill.fore_color.rgb = C["navy"]

    footer_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(5.35),
        Inches(9.2), Inches(0.2)
    )
    footer_frame = footer_tf.text_frame
    footer_frame.word_wrap = True
    footer_p = footer_frame.paragraphs[0]
    footer_p.text = "Fuente: DSFFA — MIDAGRI, SAC 2025-2026"
    footer_p.font.size = Pt(8)
    footer_p.font.italic = True
    footer_p.font.color.rgb = C["gray"]
    footer_p.alignment = PP_ALIGN.RIGHT


def _add_portada(prs, data):
    """Add cover slide with dark navy background and teal accents."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["navy"]
    background.line.fill.background()

    line_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.35),
        prs.slide_width, Inches(0.06)
    )
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = C["teal"]
    line_top.line.fill.background()

    line_bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.22),
        prs.slide_width, Inches(0.06)
    )
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = C["teal"]
    line_bottom.line.fill.background()

    _make_logo(slide, Inches(4.0), Inches(0.55), Inches(2.0), Inches(1.6))

    tf_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.3),
        Inches(9), Inches(0.55)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    text_frame.margin_top = Pt(0)
    text_frame.margin_bottom = Pt(0)
    p = text_frame.paragraphs[0]
    p.text = "SEGURO AGRÍCOLA CATASTRÓFICO"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C["white"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    tf_subtitle = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.95),
        Inches(9), Inches(0.35)
    )
    text_frame = tf_subtitle.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "SAC 2025–2026"
    p.font.size = Pt(20)
    p.font.color.rgb = C["teal"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    sep_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(3.45),
        Inches(3), Inches(0.03)
    )
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = C["teal"]
    sep_line.line.fill.background()

    # ── Índice / contenido de la presentación ──
    filtros = data.get("filtros", {})
    indice_lines = []

    # Nivel 1: estructura geográfica
    if data.get("scope") == "nacional" or (not filtros.get("deptos") and not filtros.get("provs") and not filtros.get("dists")):
        indice_lines.append("Nacional")
    if filtros.get("deptos"):
        for d in filtros["deptos"]:
            indice_lines.append(f"Departamento: {_fix_geo_name(d)}")
    if filtros.get("provs"):
        for pv in filtros["provs"]:
            indice_lines.append(f"  Provincia: {_fix_geo_name(pv)}")
    if filtros.get("dists"):
        for d in filtros["dists"][:5]:
            indice_lines.append(f"    Distrito: {_fix_geo_name(d)}")

    # Nivel 2
    if data.get("hay_nivel2"):
        indice_lines.append("")
        indice_lines.append(f"Análisis Complementario: {data.get('nivel2_label', '')}")

    tf_scope = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.55),
        Inches(7), Inches(1.2)
    )
    text_frame = tf_scope.text_frame
    text_frame.word_wrap = True

    for idx, line in enumerate(indice_lines):
        if idx == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = C["lightGray"] if line.strip().startswith(("Provincia", "Distrito", "Análisis")) else C["white"]
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(1)
        p.space_after = Pt(1)

    fecha_corte = data.get("fecha_corte", "S.F.")
    tf_footer = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.85),
        Inches(9), Inches(0.35)
    )
    text_frame = tf_footer.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Dirección de Seguro y Fomento del Financiamiento Agrario — MIDAGRI"
    p.font.size = Pt(10)
    p.font.italic = True
    p.font.color.rgb = C["lightGray"]
    p.alignment = PP_ALIGN.CENTER

    p2 = text_frame.add_paragraph()
    p2.text = f"Corte al {fecha_corte}"
    p2.font.size = Pt(9)
    p2.font.italic = True
    p2.font.color.rgb = C["lightGray"]
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(2)


def _add_cierre(prs, fecha_corte):
    """Add closing slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["navy"]
    background.line.fill.background()

    line_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.35),
        prs.slide_width, Inches(0.06)
    )
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = C["teal"]
    line_top.line.fill.background()

    line_bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.22),
        prs.slide_width, Inches(0.06)
    )
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = C["teal"]
    line_bottom.line.fill.background()

    _make_logo(slide, Inches(4.0), Inches(0.55), Inches(2.0), Inches(1.3))

    tf_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.0),
        Inches(9), Inches(0.55)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    text_frame.margin_top = Pt(0)
    text_frame.margin_bottom = Pt(0)
    p = text_frame.paragraphs[0]
    p.text = "SEGURO AGRÍCOLA CATASTRÓFICO"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = C["white"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    tf_subtitle = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.65),
        Inches(9), Inches(0.35)
    )
    text_frame = tf_subtitle.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "SAC 2025–2026"
    p.font.size = Pt(20)
    p.font.color.rgb = C["teal"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    sep_line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(3.5), Inches(3.05),
        Inches(3), Inches(0.03)
    )
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = C["teal"]
    sep_line.line.fill.background()

    tf_footer = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.4),
        Inches(9), Inches(1.0)
    )
    text_frame = tf_footer.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Dirección de Seguro y Fomento del Financiamiento Agrario"
    p.font.size = Pt(13)
    p.font.color.rgb = C["white"]
    p.alignment = PP_ALIGN.CENTER

    p2 = text_frame.add_paragraph()
    p2.text = "Ministerio de Desarrollo Agrario y Riego — MIDAGRI"
    p2.font.size = Pt(12)
    p2.font.color.rgb = C["lightGray"]
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

    p3 = text_frame.add_paragraph()
    p3.text = f"Corte al {fecha_corte}"
    p3.font.size = Pt(10)
    p3.font.italic = True
    p3.font.color.rgb = C["lightGray"]
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(4)


def _add_nacional_section(prs, section):
    """Add nacional section: indicadores clave, pipeline, tipo siniestro, top deptos."""
    m = section["metricas"]
    empresas = section.get("empresas", [])
    pipeline = section.get("pipeline", [])
    tipos = section.get("tipos", [])
    deptos = section.get("departamentos", [])

    # Slide 1: Indicadores Clave
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["cream"]
    background.line.fill.background()

    _make_logo(slide, Inches(0.4), Inches(0.35), Inches(0.6), Inches(0.65))

    tf_title = slide.shapes.add_textbox(
        Inches(1.2), Inches(0.3),
        Inches(8.3), Inches(0.55)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Resumen Nacional — Indicadores Clave"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C["navy"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    tf_sub = slide.shapes.add_textbox(
        Inches(1.2), Inches(0.85),
        Inches(8.3), Inches(0.25)
    )
    text_frame = tf_sub.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    emp_text = "Consolidado La Positiva + Rímac · Campaña 2025-2026"
    if empresas:
        emp_text = " + ".join([f"{e['empresa']} ({e['avisos']} avisos)" for e in empresas[:2]])
    p.text = emp_text
    p.font.size = Pt(10)
    p.font.color.rgb = C["gray"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.15),
        prs.slide_width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C["teal"]
    line.line.fill.background()

    kpi_configs = [
        ("Avisos Reportados", _fmt_num(m["avisos"]), "24 departamentos", C["teal"], "⚠"),
        ("Avance Evaluación", f"{_fmt_pct(m['pct_eval'])}", f"{m['cerrados']:,} cerrados de {m['avisos']:,}", C["teal"], "✔"),
        ("Indemnización", _fmt_money(m["indemnizacion"]), "Reconocida a productores", C["orange"], "●"),
        ("Avance Desembolso", f"{_fmt_pct(m['pct_desembolso'])}", f"{_fmt_money(m['desembolso'])} desembolsados", C["teal"], "✦"),
        ("Ha Indemnizadas", f"{m['ha_indemnizadas']:,.1f}", "Solo con evaluación cerrada", C["teal"], "🌿"),
        ("Productores", _fmt_num(m["productores"]), "Beneficiados con indemnización", C["sage"], "👥"),
    ]

    if empresas:
        for i, emp in enumerate(empresas[:2]):
            emp_label = f"{emp['empresa']}"
            emp_val = f"{_fmt_num(emp['avisos'])}"
            emp_detail = f"{_fmt_pct(emp['pct_eval'])} cerr. · {_fmt_money(emp['indemnizacion'])}"
            kpi_configs.append((emp_label, emp_val, emp_detail, C["teal"], f"{i+1}"))

    for i, (label, value, sublabel, color, icon) in enumerate(kpi_configs[:8]):
        col = i % 4
        row = i // 4
        left = Inches(0.35 + col * 2.35)
        top = Inches(1.65 + row * 2.0)
        _add_kpi_card(slide, left, top, Inches(2.15), Inches(1.85), label, value, sublabel, color, icon)

    # Slide 2: Pipeline de Procesos
    if pipeline:
        _add_pipeline_slide(prs, pipeline, {}, m)

    # Slide 3: Tipo de Siniestro
    if tipos:
        _add_tipo_siniestro_slide(prs, tipos)

    # Slide 4: Top Departamentos
    if deptos:
        _add_top_deptos_chart(prs, deptos)


def _add_departamental_section(prs, section, fecha_corte="S.F."):
    """Add departamental section: separator + metrics + resumen ejecutivo."""
    name = section.get("name", "Departamento")
    m = section["metricas"]
    tipos = section.get("tipos", [])
    provs = section.get("provincias", [])
    pipeline = section.get("pipeline", [])
    insights = section.get("insights", [])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["navy"]
    background.line.fill.background()

    line_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.35),
        prs.slide_width, Inches(0.06)
    )
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = C["teal"]
    line_top.line.fill.background()

    line_bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.22),
        prs.slide_width, Inches(0.06)
    )
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = C["teal"]
    line_bottom.line.fill.background()

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.15), Inches(0.8),
        Inches(1.7), Inches(1.7)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = C["teal"]
    circle.line.fill.background()

    circle_tf = slide.shapes.add_textbox(
        Inches(4.15), Inches(0.8),
        Inches(1.7), Inches(1.7)
    )
    circle_frame = circle_tf.text_frame
    circle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    circle_p = circle_frame.paragraphs[0]
    circle_p.text = "🌾"
    circle_p.font.size = Pt(60)
    circle_p.alignment = PP_ALIGN.CENTER

    # Etiqueta de nivel geográfico sobre el nombre
    tf_level = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.35),
        Inches(9), Inches(0.35)
    )
    level_frame = tf_level.text_frame
    level_frame.word_wrap = True
    lp = level_frame.paragraphs[0]
    lp.text = "DEPARTAMENTO"
    lp.font.size = Pt(14)
    lp.font.color.rgb = C["teal"]
    lp.font.name = "Calibri"
    lp.alignment = PP_ALIGN.CENTER

    tf_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.7),
        Inches(9), Inches(0.8)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = _fix_geo_name(name)
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C["white"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    p2 = text_frame.add_paragraph()
    p2.text = f"{m['avisos']:,} avisos — {_fmt_money(m['indemnizacion'])} indemnización"
    p2.font.size = Pt(13)
    p2.font.color.rgb = C["teal"]
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

    # ── Slide 2: KPIs + Tablas ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["cream"]
    background.line.fill.background()

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(0.35),
        Inches(2.2), Inches(0.35)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = C["white"]
    badge.line.color.rgb = C["orange"]
    badge.line.width = Pt(1.5)

    badge_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.35),
        Inches(2.2), Inches(0.35)
    )
    badge_frame = badge_tf.text_frame
    badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    badge_p = badge_frame.paragraphs[0]
    badge_p.text = f"DEPTO: {_fix_geo_upper(name)}"
    badge_p.font.size = Pt(11)
    badge_p.font.bold = True
    badge_p.font.color.rgb = C["orange"]
    badge_p.alignment = PP_ALIGN.CENTER

    tf_title = slide.shapes.add_textbox(
        Inches(2.8), Inches(0.28),
        Inches(6.8), Inches(0.5)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Análisis Departamental — Indicadores Clave"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C["navy"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.85),
        prs.slide_width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C["teal"]
    line.line.fill.background()

    kpi_configs = [
        ("Avisos Reportados", _fmt_num(m["avisos"]), "", C["amber"], "⚠"),
        ("Avance Evaluación", f"{_fmt_pct(m['pct_eval'])}", f"{m['cerrados']:,} cerrados", C["teal"], "✔"),
        ("Indemnización", _fmt_money(m["indemnizacion"]), "", C["orange"], "●"),
        ("Avance Desembolso", f"{_fmt_pct(m['pct_desembolso'])}", _fmt_money(m["desembolso"]), C["teal"], "✦"),
    ]

    for i, (label, value, sublabel, color, icon) in enumerate(kpi_configs):
        left = Inches(0.3 + i * 2.35)
        top = Inches(1.25)
        _add_kpi_card(slide, left, top, Inches(2.15), Inches(1.85), label, value, sublabel, color, icon)

    if tipos:
        headers = ["Tipo Siniestro", "Avisos", "Indem. (S/)"]
        rows = []
        for t in tipos[:8]:
            rows.append([
                t["tipo"],
                _fmt_num(t["avisos"]),
                _fmt_money(t.get("indem", 0)),
            ])
        col_widths = [Inches(2.0), Inches(1.1), Inches(1.2)]
        _add_styled_table(slide, headers, rows, left=Inches(0.3), top=Inches(3.25), col_widths=col_widths, max_rows=5)

    if provs:
        headers_p = ["Provincia", "Avisos", "Indem. (S/)"]
        rows_p = []
        for p in provs[:5]:
            rows_p.append([
                p["name"],
                _fmt_num(p["avisos"]),
                _fmt_money(p.get("indem", 0)),
            ])
        col_widths_p = [Inches(2.0), Inches(1.1), Inches(1.2)]
        _add_styled_table(slide, headers_p, rows_p, left=Inches(5.0), top=Inches(3.25), col_widths=col_widths_p, max_rows=5)

    # Slide 3: Resumen Ejecutivo
    resumen_text = _generar_resumen_texto(section, "departamental")
    _add_resumen_ejecutivo(prs, name, resumen_text, "Departamental", fecha_corte)


def _add_provincial_section(prs, section, fecha_corte="S.F."):
    """Add provincial section: separator + metrics + tables + resumen ejecutivo."""
    name = section.get("name", "Provincia")
    depto = section.get("depto", "")
    m = section["metricas"]
    dists = section.get("distritos", [])
    tipos = section.get("tipos", [])

    # ── Slide 1: Separador provincial ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["forest"]
    background.line.fill.background()

    line_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.35),
        prs.slide_width, Inches(0.06)
    )
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = C["sage"]
    line_top.line.fill.background()

    line_bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.22),
        prs.slide_width, Inches(0.06)
    )
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = C["sage"]
    line_bottom.line.fill.background()

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.15), Inches(0.8),
        Inches(1.7), Inches(1.7)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = C["sage"]
    circle.line.fill.background()

    circle_tf = slide.shapes.add_textbox(
        Inches(4.15), Inches(0.8),
        Inches(1.7), Inches(1.7)
    )
    circle_frame = circle_tf.text_frame
    circle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    circle_p = circle_frame.paragraphs[0]
    circle_p.text = "🏡"
    circle_p.font.size = Pt(60)
    circle_p.alignment = PP_ALIGN.CENTER

    # Etiqueta de nivel geográfico
    tf_level = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.35),
        Inches(9), Inches(0.35)
    )
    level_frame = tf_level.text_frame
    level_frame.word_wrap = True
    lp = level_frame.paragraphs[0]
    lp.text = "PROVINCIA"
    lp.font.size = Pt(14)
    lp.font.color.rgb = C["sage"]
    lp.font.name = "Calibri"
    lp.alignment = PP_ALIGN.CENTER

    tf_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.7),
        Inches(9), Inches(0.8)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = _fix_geo_name(name)
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C["white"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    p2 = text_frame.add_paragraph()
    p2.text = f"{_fix_geo_name(depto)} — {m['avisos']:,} avisos — {_fmt_money(m['indemnizacion'])} indemnización" if depto else "Provincia"
    p2.font.size = Pt(13)
    p2.font.color.rgb = C["sage"]
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

    # ── Slide 2: KPIs + Tablas (tipo siniestro y distritos) ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["cream"]
    background.line.fill.background()

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(0.35),
        Inches(2.2), Inches(0.35)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = C["white"]
    badge.line.color.rgb = C["sage"]
    badge.line.width = Pt(1.5)

    badge_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.35),
        Inches(2.2), Inches(0.35)
    )
    badge_frame = badge_tf.text_frame
    badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    badge_p = badge_frame.paragraphs[0]
    badge_p.text = f"PROV: {_fix_geo_upper(name)}"
    badge_p.font.size = Pt(11)
    badge_p.font.bold = True
    badge_p.font.color.rgb = C["forest"]
    badge_p.alignment = PP_ALIGN.CENTER

    tf_title = slide.shapes.add_textbox(
        Inches(2.8), Inches(0.28),
        Inches(6.8), Inches(0.5)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Análisis Provincial — Indicadores Clave"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C["forest"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.85),
        prs.slide_width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C["sage"]
    line.line.fill.background()

    kpi_configs = [
        ("Avisos Reportados", _fmt_num(m["avisos"]), "", C["amber"], "⚠"),
        ("Avance Evaluación", f"{_fmt_pct(m['pct_eval'])}", f"{m['cerrados']:,} cerrados", C["teal"], "✔"),
        ("Indemnización", _fmt_money(m["indemnizacion"]), "", C["orange"], "●"),
        ("Avance Desembolso", f"{_fmt_pct(m['pct_desembolso'])}", _fmt_money(m["desembolso"]), C["teal"], "✦"),
    ]

    for i, (label, value, sublabel, color, icon) in enumerate(kpi_configs):
        left = Inches(0.3 + i * 2.35)
        top = Inches(1.25)
        _add_kpi_card(slide, left, top, Inches(2.15), Inches(1.85), label, value, sublabel, color, icon)

    if tipos:
        headers = ["Tipo Siniestro", "Avisos", "Indem. (S/)"]
        rows = []
        for t in tipos[:8]:
            rows.append([
                t["tipo"],
                _fmt_num(t["avisos"]),
                _fmt_money(t.get("indem", 0)),
            ])
        col_widths = [Inches(2.0), Inches(1.1), Inches(1.2)]
        _add_styled_table(slide, headers, rows, left=Inches(0.3), top=Inches(3.25), col_widths=col_widths, max_rows=5)

    if dists:
        headers_d = ["Distrito", "Avisos", "Indem. (S/)"]
        rows_d = []
        for d in dists[:5]:
            rows_d.append([
                d["name"],
                _fmt_num(d["avisos"]),
                _fmt_money(d.get("indem", 0)),
            ])
        col_widths_d = [Inches(2.0), Inches(1.1), Inches(1.2)]
        _add_styled_table(slide, headers_d, rows_d, left=Inches(5.0), top=Inches(3.25), col_widths=col_widths_d, max_rows=5)

    # ── Slide 3: Resumen Ejecutivo ──
    resumen_text = _generar_resumen_texto(section, "provincial")
    _add_resumen_ejecutivo(prs, name, resumen_text, "Provincial", fecha_corte)


def _add_distrital_section(prs, section, fecha_corte="S.F."):
    """Add distrital section: separator + metrics + tables + resumen ejecutivo."""
    name = section.get("name", "Distrito")
    prov = section.get("prov", "")
    depto = section.get("depto", "")
    m = section["metricas"]
    tipos = section.get("tipos", [])

    # ── Slide 1: Separador distrital ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["navy"]
    background.line.fill.background()

    line_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.35),
        prs.slide_width, Inches(0.06)
    )
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = C["teal"]
    line_top.line.fill.background()

    line_bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.22),
        prs.slide_width, Inches(0.06)
    )
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = C["teal"]
    line_bottom.line.fill.background()

    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.15), Inches(0.8),
        Inches(1.7), Inches(1.7)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = C["teal"]
    circle.line.fill.background()

    circle_tf = slide.shapes.add_textbox(
        Inches(4.15), Inches(0.8),
        Inches(1.7), Inches(1.7)
    )
    circle_frame = circle_tf.text_frame
    circle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    circle_p = circle_frame.paragraphs[0]
    circle_p.text = "📍"
    circle_p.font.size = Pt(60)
    circle_p.alignment = PP_ALIGN.CENTER

    # Etiqueta de nivel geográfico
    tf_level = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.35),
        Inches(9), Inches(0.35)
    )
    level_frame = tf_level.text_frame
    level_frame.word_wrap = True
    lp = level_frame.paragraphs[0]
    lp.text = "DISTRITO"
    lp.font.size = Pt(14)
    lp.font.color.rgb = C["teal"]
    lp.font.name = "Calibri"
    lp.alignment = PP_ALIGN.CENTER

    tf_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.7),
        Inches(9), Inches(0.8)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = _fix_geo_name(name)
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = C["white"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    subtitle_parts = []
    if prov:
        subtitle_parts.append(_fix_geo_name(prov))
    if depto:
        subtitle_parts.append(_fix_geo_name(depto))
    subtitle = " — ".join(subtitle_parts) if subtitle_parts else "Distrito"
    p2 = text_frame.add_paragraph()
    p2.text = f"{subtitle} — {m['avisos']:,} avisos — {_fmt_money(m['indemnizacion'])} indemnización"
    p2.font.size = Pt(13)
    p2.font.color.rgb = C["teal"]
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(4)

    # ── Slide 2: KPIs + Tabla de tipos de siniestro ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["cream"]
    background.line.fill.background()

    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(0.35),
        Inches(2.2), Inches(0.35)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = C["white"]
    badge.line.color.rgb = C["teal"]
    badge.line.width = Pt(1.5)

    badge_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.35),
        Inches(2.2), Inches(0.35)
    )
    badge_frame = badge_tf.text_frame
    badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    badge_p = badge_frame.paragraphs[0]
    badge_p.text = f"DIST: {_fix_geo_upper(name)}"
    badge_p.font.size = Pt(11)
    badge_p.font.bold = True
    badge_p.font.color.rgb = C["navy"]
    badge_p.alignment = PP_ALIGN.CENTER

    tf_title = slide.shapes.add_textbox(
        Inches(2.8), Inches(0.28),
        Inches(6.8), Inches(0.5)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Análisis Distrital — Indicadores Clave"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C["navy"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.85),
        prs.slide_width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C["teal"]
    line.line.fill.background()

    kpi_configs = [
        ("Avisos Reportados", _fmt_num(m["avisos"]), "", C["amber"], "⚠"),
        ("Avance Evaluación", f"{_fmt_pct(m['pct_eval'])}", f"{m['cerrados']:,} cerrados", C["teal"], "✔"),
        ("Indemnización", _fmt_money(m["indemnizacion"]), "", C["orange"], "●"),
        ("Avance Desembolso", f"{_fmt_pct(m['pct_desembolso'])}", _fmt_money(m["desembolso"]), C["teal"], "✦"),
    ]

    for i, (label, value, sublabel, color, icon) in enumerate(kpi_configs):
        left = Inches(0.3 + i * 2.35)
        top = Inches(1.25)
        _add_kpi_card(slide, left, top, Inches(2.15), Inches(1.85), label, value, sublabel, color, icon)

    if tipos:
        headers = ["Tipo Siniestro", "Avisos", "Indem. (S/)"]
        rows = []
        for t in tipos[:8]:
            rows.append([
                t["tipo"],
                _fmt_num(t["avisos"]),
                _fmt_money(t.get("indem", 0)),
            ])
        col_widths = [Inches(2.0), Inches(1.1), Inches(1.2)]
        _add_styled_table(slide, headers, rows, left=Inches(0.3), top=Inches(3.25), col_widths=col_widths, max_rows=5)

    # ── Slide 3: Resumen Ejecutivo ──
    resumen_text = _generar_resumen_texto(section, "distrital")
    _add_resumen_ejecutivo(prs, name, resumen_text, "Distrital", fecha_corte)


# ══════════════════════════════════════════════════════════════════
# NIVEL 2 — SLIDE GENERATORS (complementary filtered sections)
# ══════════════════════════════════════════════════════════════════

def _add_nivel2_separator(prs, nivel2_label, fecha_corte):
    """Add dark separator slide that introduces Level 2 (complementary analysis)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["forest"]
    background.line.fill.background()

    line_top = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0.35),
        prs.slide_width, Inches(0.06)
    )
    line_top.fill.solid()
    line_top.fill.fore_color.rgb = C["gold"]
    line_top.line.fill.background()

    line_bottom = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(5.22),
        prs.slide_width, Inches(0.06)
    )
    line_bottom.fill.solid()
    line_bottom.fill.fore_color.rgb = C["gold"]
    line_bottom.line.fill.background()

    # Icon circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.15), Inches(0.7),
        Inches(1.7), Inches(1.7)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = C["gold"]
    circle.line.fill.background()

    circle_tf = slide.shapes.add_textbox(
        Inches(4.15), Inches(0.7),
        Inches(1.7), Inches(1.7)
    )
    circle_frame = circle_tf.text_frame
    circle_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    circle_p = circle_frame.paragraphs[0]
    circle_p.text = "🔍"
    circle_p.font.size = Pt(60)
    circle_p.alignment = PP_ALIGN.CENTER

    tf_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.6),
        Inches(9), Inches(0.6)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "ANÁLISIS COMPLEMENTARIO"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C["white"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    tf_sub = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.25),
        Inches(9), Inches(0.5)
    )
    text_frame = tf_sub.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = nivel2_label
    p.font.size = Pt(16)
    p.font.color.rgb = C["gold"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER

    tf_desc = slide.shapes.add_textbox(
        Inches(1.0), Inches(3.9),
        Inches(8), Inches(0.5)
    )
    text_frame = tf_desc.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Las siguientes secciones muestran el mismo alcance geográfico filtrado por las características seleccionadas."
    p.font.size = Pt(11)
    p.font.italic = True
    p.font.color.rgb = C["lightGray"]
    p.alignment = PP_ALIGN.CENTER

    tf_footer = slide.shapes.add_textbox(
        Inches(0.5), Inches(4.85),
        Inches(9), Inches(0.3)
    )
    text_frame = tf_footer.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = f"Corte al {fecha_corte}"
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = C["lightGray"]
    p.alignment = PP_ALIGN.CENTER


def _add_nivel2_nacional(prs, section, nivel2_label):
    """Add Level 2 nacional slide: KPI cards + tipo siniestro for filtered data."""
    m = section["metricas"]
    tipos = section.get("tipos", [])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["cream"]
    background.line.fill.background()

    # Badge for Level 2
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(0.3),
        Inches(2.0), Inches(0.3)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = C["gold"]
    badge.line.fill.background()

    badge_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.3),
        Inches(2.0), Inches(0.3)
    )
    badge_frame = badge_tf.text_frame
    badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    badge_p = badge_frame.paragraphs[0]
    badge_p.text = "COMPLEMENTARIO"
    badge_p.font.size = Pt(9)
    badge_p.font.bold = True
    badge_p.font.color.rgb = C["white"]
    badge_p.alignment = PP_ALIGN.CENTER

    tf_title = slide.shapes.add_textbox(
        Inches(2.6), Inches(0.25),
        Inches(7.0), Inches(0.45)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Nacional — Análisis Filtrado"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C["navy"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    tf_sub = slide.shapes.add_textbox(
        Inches(2.6), Inches(0.72),
        Inches(7.0), Inches(0.25)
    )
    text_frame = tf_sub.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = nivel2_label
    p.font.size = Pt(10)
    p.font.color.rgb = C["gray"]
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.05),
        prs.slide_width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C["gold"]
    line.line.fill.background()

    # 4 KPI cards
    kpi_configs = [
        ("Avisos Filtrados", _fmt_num(m["avisos"]), "", C["gold"], "⚠"),
        ("Avance Evaluación", _fmt_pct(m["pct_eval"]), f"{m['cerrados']:,} cerrados", C["teal"], "✔"),
        ("Indemnización", _fmt_money(m["indemnizacion"]), "Reconocida", C["orange"], "●"),
        ("Avance Desembolso", _fmt_pct(m["pct_desembolso"]), _fmt_money(m["desembolso"]), C["teal"], "✦"),
    ]

    for i, (label, value, sublabel, color, icon) in enumerate(kpi_configs):
        left = Inches(0.35 + i * 2.35)
        top = Inches(1.35)
        _add_kpi_card(slide, left, top, Inches(2.15), Inches(1.85), label, value, sublabel, color, icon)

    # Tabla de tipos — máx 4 filas para que quepa dentro de la lámina
    if tipos:
        headers = ["Tipo Siniestro", "Avisos", "Indem. (S/)"]
        rows = []
        for t in tipos[:4]:
            rows.append([t["tipo"], _fmt_num(t["avisos"]), _fmt_money(t.get("indem", 0))])
        _add_styled_table(slide, headers, rows, left=Inches(0.5), top=Inches(3.45),
                         col_widths=[Inches(3.5), Inches(1.5), Inches(2.0)], max_rows=4)


def _add_nivel2_departamental(prs, section, nivel2_label, fecha_corte):
    """Add Level 2 departamental slide: KPI cards + tables for filtered data."""
    name = section.get("name", "Departamento")
    m = section["metricas"]
    tipos = section.get("tipos", [])
    provs = section.get("provincias", [])

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = C["cream"]
    background.line.fill.background()

    # Badge: dept name + complementary tag
    badge = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.4), Inches(0.3),
        Inches(1.5), Inches(0.3)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = C["white"]
    badge.line.color.rgb = C["gold"]
    badge.line.width = Pt(1.5)

    badge_tf = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.3),
        Inches(1.5), Inches(0.3)
    )
    badge_frame = badge_tf.text_frame
    badge_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    badge_p = badge_frame.paragraphs[0]
    badge_p.text = _fix_geo_upper(name)
    badge_p.font.size = Pt(10)
    badge_p.font.bold = True
    badge_p.font.color.rgb = C["gold"]
    badge_p.alignment = PP_ALIGN.CENTER

    # Complementario tag
    tag = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(2.05), Inches(0.3),
        Inches(1.8), Inches(0.3)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = C["gold"]
    tag.line.fill.background()

    tag_tf = slide.shapes.add_textbox(
        Inches(2.05), Inches(0.3),
        Inches(1.8), Inches(0.3)
    )
    tag_frame = tag_tf.text_frame
    tag_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    tag_p = tag_frame.paragraphs[0]
    tag_p.text = "COMPLEMENTARIO"
    tag_p.font.size = Pt(8)
    tag_p.font.bold = True
    tag_p.font.color.rgb = C["white"]
    tag_p.alignment = PP_ALIGN.CENTER

    tf_title = slide.shapes.add_textbox(
        Inches(4.0), Inches(0.25),
        Inches(5.6), Inches(0.45)
    )
    text_frame = tf_title.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = "Análisis Filtrado"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = C["navy"]
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.LEFT

    tf_sub = slide.shapes.add_textbox(
        Inches(4.0), Inches(0.72),
        Inches(5.6), Inches(0.25)
    )
    text_frame = tf_sub.text_frame
    text_frame.word_wrap = True
    p = text_frame.paragraphs[0]
    p.text = nivel2_label
    p.font.size = Pt(9)
    p.font.color.rgb = C["gray"]
    p.alignment = PP_ALIGN.LEFT

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(1.05),
        prs.slide_width, Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C["gold"]
    line.line.fill.background()

    # 4 KPI cards
    kpi_configs = [
        ("Avisos Filtrados", _fmt_num(m["avisos"]), "", C["gold"], "⚠"),
        ("Avance Evaluación", _fmt_pct(m["pct_eval"]), f"{m['cerrados']:,} cerrados", C["teal"], "✔"),
        ("Indemnización", _fmt_money(m["indemnizacion"]), "", C["orange"], "●"),
        ("Avance Desembolso", _fmt_pct(m["pct_desembolso"]), _fmt_money(m["desembolso"]), C["teal"], "✦"),
    ]

    for i, (label, value, sublabel, color, icon) in enumerate(kpi_configs):
        left = Inches(0.3 + i * 2.35)
        top = Inches(1.25)
        _add_kpi_card(slide, left, top, Inches(2.15), Inches(1.85), label, value, sublabel, color, icon)

    # Tables: tipos + provincias (same layout as departamental)
    if tipos:
        headers = ["Tipo Siniestro", "Avisos", "Indem. (S/)"]
        rows = [[t["tipo"], _fmt_num(t["avisos"]), _fmt_money(t.get("indem", 0))] for t in tipos[:5]]
        _add_styled_table(slide, headers, rows, left=Inches(0.3), top=Inches(3.25),
                         col_widths=[Inches(2.0), Inches(1.1), Inches(1.2)], max_rows=5)

    if provs:
        headers_p = ["Provincia", "Avisos", "Indem. (S/)"]
        rows_p = [[p["name"], _fmt_num(p["avisos"]), _fmt_money(p.get("indem", 0))] for p in provs[:5]]
        _add_styled_table(slide, headers_p, rows_p, left=Inches(5.0), top=Inches(3.25),
                         col_widths=[Inches(2.0), Inches(1.1), Inches(1.2)], max_rows=5)

    # Resumen ejecutivo for Level 2
    resumen_text = _generar_resumen_texto(section, "departamental")
    _add_resumen_ejecutivo(prs, f"{name} — Filtrado", resumen_text, "Complementario", fecha_corte)


# ══════════════════════════════════════════════════════════════════
# MAIN GENERATION FUNCTION
# ══════════════════════════════════════════════════════════════════

def generar_ppt_dinamico(df, filtros, fecha_corte):
    """
    Genera una presentación PPT dinámica con python-pptx.
    Modelo de DOS NIVELES ACUMULATIVOS:
      Nivel 1: Base geográfica (Nacional → Depto → Prov → Dist)
      Nivel 2: Análisis complementario con filtros adicionales
               (tipos de siniestro, fechas de ocurrencia, etc.)

    Args:
        df: DataFrame consolidado (datos["midagri"])
        filtros: dict con selecciones del usuario
        fecha_corte: string con fecha de corte (ej: "14/03/2026")

    Returns:
        bytes del archivo .pptx
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    prs.author = "DSFFA — MIDAGRI"
    prs.title = "SAC 2025-2026 — Presentación Dinámica"

    data = _prepare_data(df, filtros, fecha_corte)

    # ── Portada ──
    _add_portada(prs, data)

    # ── NIVEL 1: Secciones geográficas (base) ──
    for section in data.get("nivel1_sections", []):
        section_type = section.get("type", "")

        if section_type == "nacional":
            _add_nacional_section(prs, section)
        elif section_type == "departamental":
            _add_departamental_section(prs, section, fecha_corte)
        elif section_type == "provincial":
            _add_provincial_section(prs, section, fecha_corte)
        elif section_type == "distrital":
            _add_distrital_section(prs, section, fecha_corte)

    # ── NIVEL 2: Secciones complementarias (filtradas) ──
    nivel2_sections = data.get("nivel2_sections", [])
    if nivel2_sections and data.get("hay_nivel2"):
        nivel2_label = data.get("nivel2_label", "")

        # Separador de Nivel 2
        _add_nivel2_separator(prs, nivel2_label, fecha_corte)

        for section in nivel2_sections:
            section_type = section.get("type", "")

            if section_type == "nacional":
                _add_nivel2_nacional(prs, section, nivel2_label)
            elif section_type == "departamental":
                _add_nivel2_departamental(prs, section, nivel2_label, fecha_corte)
            elif section_type == "provincial":
                _add_provincial_section(prs, section, fecha_corte)

    # ── Cierre ──
    _add_cierre(prs, fecha_corte)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.getvalue()

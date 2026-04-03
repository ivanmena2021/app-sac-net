"""
gen_ppt_historico.py — PPT de análisis histórico de siniestralidad por departamento
===================================================================================
Genera una presentación de 12 slides con:
  - Siniestralidad = Indemnización / Prima Neta (%) por campaña
  - Datos históricos de 5 campañas (JSONs estáticos)
  - Campaña actual (datos dinámicos de datos["midagri"])
  - Primas netas históricas (Excel estático)

Fuentes:
  - static_data/resumen_departamental.json (montos por campaña×depto)
  - static_data/Primas_Totales_SAC_2020-2026.xlsx (primas netas)
  - datos["midagri"] (campaña actual en tiempo real)
  - datos["materia"] (prima neta actual)

Metodología: static_data/METODOLOGIA_DATOS.md
"""
import io
import json
import os
from datetime import datetime

import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.chart.data import CategoryChartData

# ═══════════════════════════════════════════════════════════════
# COLORES MIDAGRI
# ═══════════════════════════════════════════════════════════════

VERDE = RGBColor(0x40, 0x8B, 0x14)
TEAL = RGBColor(0x10, 0xA9, 0xA7)
GRIS = RGBColor(0x3F, 0x3F, 0x3F)
DORADO = RGBColor(0xFF, 0xC0, 0x00)
BLANCO = RGBColor(0xFF, 0xFF, 0xFF)
CREMA = RGBColor(0xF2, 0xF2, 0xF2)
ROJO = RGBColor(0xE7, 0x4C, 0x3C)
NARANJA = RGBColor(0xF3, 0x9C, 0x12)
AZUL = RGBColor(0x29, 0x80, 0xB9)
FONT = "Calibri"

CAMPANAS_HIST = ["2020-2021", "2021-2022", "2022-2023", "2023-2024", "2024-2025"]

# Cargar datos estáticos
_STATIC_DIR = os.path.join(os.path.dirname(__file__), "static_data")
_RESUMEN_PATH = os.path.join(_STATIC_DIR, "resumen_departamental.json")

_RESUMEN_DEPT = {}
if os.path.exists(_RESUMEN_PATH):
    with open(_RESUMEN_PATH, "r", encoding="utf-8") as f:
        _RESUMEN_DEPT = json.load(f)


# ═══════════════════════════════════════════════════════════════
# HELPERS
# ═══════════════════════════════════════════════════════════════

def _safe_chart_val(val, default=0):
    """Sanitize numeric value for chart data: NaN/inf/None -> default."""
    if val is None:
        return default
    try:
        f = float(val)
        if np.isnan(f) or np.isinf(f):
            return default
        return f
    except (ValueError, TypeError):
        return default


def _bg(slide, prs, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid(); bg.fill.fore_color.rgb = color; bg.line.fill.background()


def _text(slide, x, y, w, h, txt, size=14, bold=False, color=GRIS, align=PP_ALIGN.LEFT, font=FONT):
    tf = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = tf.text_frame; frame.word_wrap = True
    p = frame.paragraphs[0]; p.text = str(txt)
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    p.font.name = font; p.alignment = align
    return frame


def _lines(slide, prs):
    for y in [0.4, 6.8]:
        ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(y), prs.slide_width, Inches(0.06))
        ln.fill.solid(); ln.fill.fore_color.rgb = VERDE; ln.line.fill.background()


def _kpi_card(slide, x, y, w, h, label, value, sub, accent):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid(); card.fill.fore_color.rgb = BLANCO; card.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    _text(slide, x + 0.18, y + 0.08, w - 0.3, 0.25, label, 10, True, RGBColor(0x64, 0x74, 0x8B))
    _text(slide, x + 0.18, y + 0.35, w - 0.3, 0.55, value, 24, True, GRIS)
    if sub:
        _text(slide, x + 0.18, y + h - 0.35, w - 0.3, 0.25, sub, 9, False, RGBColor(0x94, 0xA3, 0xB8))


def _table_header(slide, x, y, widths, labels, color):
    for j, (lbl, w) in enumerate(zip(labels, widths)):
        cx = x + sum(widths[:j])
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(y), Inches(w), Inches(0.35))
        cell.fill.solid(); cell.fill.fore_color.rgb = color; cell.line.fill.background()
        _text(slide, cx + 0.08, y, w - 0.16, 0.35, lbl, 10, True, BLANCO,
              PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)


def _table_row(slide, x, y, widths, values, row_idx):
    bg_c = BLANCO if row_idx % 2 == 0 else CREMA
    for j, (val, w) in enumerate(zip(values, widths)):
        cx = x + sum(widths[:j])
        cell = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(y), Inches(w), Inches(0.38))
        cell.fill.solid(); cell.fill.fore_color.rgb = bg_c; cell.line.fill.background()
        _text(slide, cx + 0.08, y + 0.02, w - 0.16, 0.34, str(val), 11, False, GRIS,
              PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER)


def _obs_box(slide, x, y, w, h, text):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9); box.line.fill.background()
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = VERDE; bar.line.fill.background()
    _text(slide, x + 0.2, y + 0.08, w - 0.4, h - 0.16, text, 12, False, GRIS)


def _fmt(n):
    if n >= 1_000_000:
        return f"S/ {n/1_000_000:,.2f}M"
    if n >= 1_000:
        return f"S/ {n/1_000:,.0f}K"
    return f"S/ {n:,.0f}"


def _sin_color(pct):
    if pct > 70: return ROJO
    if pct > 50: return NARANJA
    return VERDE


def _normalize_dept(name):
    if not name:
        return ""
    n = str(name).strip().upper()
    for a, p in [("Á","A"),("É","E"),("Í","I"),("Ó","O"),("Ú","U"),("Ñ","N"),
                 ("á","A"),("é","E"),("í","I"),("ó","O"),("ú","U"),("ñ","N")]:
        n = n.replace(a, p)
    return n


# ═══════════════════════════════════════════════════════════════
# EXTRACCIÓN DE DATOS DE LA CAMPAÑA ACTUAL
# ═══════════════════════════════════════════════════════════════

def _extract_current_campaign(datos, depto):
    """Extrae métricas de la campaña actual desde datos['midagri'] (tiempo real)."""
    df = datos.get("midagri")
    if df is None or df.empty:
        return None

    depto_norm = _normalize_dept(depto)
    dept_col = "DEPARTAMENTO"
    if dept_col not in df.columns:
        return None

    puno = df[df[dept_col].astype(str).str.strip().str.upper().apply(_normalize_dept) == depto_norm].copy()
    if puno.empty:
        return None

    # Indemnización
    indemn = pd.to_numeric(puno.get("INDEMNIZACION", 0), errors="coerce").fillna(0)
    monto = float(indemn.sum())
    sup_ind = float(pd.to_numeric(puno.get("SUP_INDEMNIZADA", 0), errors="coerce").fillna(0).sum())
    desemb = float(pd.to_numeric(puno.get("MONTO_DESEMBOLSADO", 0), errors="coerce").fillna(0).sum())

    # Dictamen
    dictamen = puno.get("DICTAMEN", pd.Series(dtype=str)).astype(str).str.strip().str.upper()
    ind_mask = dictamen.str.contains("INDEMNIZABLE", na=False) & ~dictamen.str.contains("NO INDEMNIZABLE", na=False)
    no_ind = dictamen.str.contains("NO INDEMNIZABLE", na=False).sum()
    pendientes = len(puno) - ind_mask.sum() - no_ind

    # Prima neta del departamento (de materia asegurada actual)
    materia = datos.get("materia")
    prima_neta = 0
    if materia is not None and not materia.empty and "DEPARTAMENTO" in materia.columns:
        mat_dept = materia[materia["DEPARTAMENTO"].astype(str).str.strip().str.upper().apply(_normalize_dept) == depto_norm]
        if not mat_dept.empty and "PRIMA_NETA" in mat_dept.columns:
            prima_neta = float(pd.to_numeric(mat_dept["PRIMA_NETA"].iloc[0], errors="coerce") or 0)

    # Top cultivos
    cult_col = "TIPO_CULTIVO" if "TIPO_CULTIVO" in puno.columns else None
    top_cult = []
    if cult_col:
        top_cult = [[str(k), int(v)] for k, v in puno[cult_col].value_counts().head(5).items()]

    # Top siniestros
    sin_col = "TIPO_SINIESTRO" if "TIPO_SINIESTRO" in puno.columns else None
    top_sin = []
    if sin_col:
        top_sin = [[str(k), int(v)] for k, v in puno[sin_col].value_counts().head(5).items()]

    siniestralidad = round(100 * monto / prima_neta, 1) if prima_neta > 0 else 0

    return {
        "avisos": len(puno),
        "indemnizados": int(ind_mask.sum()),
        "no_indemnizables": int(no_ind),
        "pendientes": int(pendientes),
        "monto_indemnizado": monto,
        "ha_indemnizadas": sup_ind,
        "monto_desembolsado": desemb,
        "prima_neta": prima_neta,
        "siniestralidad": siniestralidad,
        "provincias": int(puno["PROVINCIA"].nunique()) if "PROVINCIA" in puno.columns else 0,
        "distritos": int(puno["DISTRITO"].nunique()) if "DISTRITO" in puno.columns else 0,
        "top_cultivos": top_cult,
        "top_siniestros": top_sin,
        "fecha_corte": datos.get("fecha_corte", "S.F."),
    }


# ═══════════════════════════════════════════════════════════════
# GENERADOR PRINCIPAL
# ═══════════════════════════════════════════════════════════════

def generar_ppt_historico(depto, datos, primas_hist):
    """Genera PPT histórica de siniestralidad para un departamento.

    Args:
        depto: nombre del departamento
        datos: dict de la app (con midagri, materia, fecha_corte)
        primas_hist: dict {campaña: {depto: prima_neta}} de load_primas_historicas()

    Returns:
        bytes del archivo .pptx
    """
    depto_norm = _normalize_dept(depto)
    depto_display = depto.strip().title()

    # Datos históricos del JSON
    hist_dept = _RESUMEN_DEPT.get("por_campana", {}).get(depto_norm, {})
    top_provs = _RESUMEN_DEPT.get("top_provincias", {}).get(depto_norm, [])
    estacionalidad = _RESUMEN_DEPT.get("estacionalidad", {}).get(depto_norm, {})

    # Construir arrays por campaña
    campanas_data = []
    for camp in CAMPANAS_HIST:
        h = hist_dept.get(camp, {})
        prima = primas_hist.get(camp, {}).get(depto_norm, 0)
        monto = h.get("monto_indemnizado", 0)
        sin_pct = round(100 * monto / prima, 1) if prima > 0 else 0
        campanas_data.append({
            "campana": camp,
            "avisos": h.get("avisos", 0),
            "indemnizados": h.get("indemnizados", 0),
            "monto": monto,
            "ha_ind": h.get("ha_indemnizadas", 0),
            "desembolso": h.get("monto_desembolsado", 0),
            "prima_neta": prima,
            "siniestralidad": sin_pct,
            "provincias": h.get("provincias", 0),
            "distritos": h.get("distritos", 0),
            "top_cultivos": h.get("top_cultivos", []),
            "top_siniestros": h.get("top_siniestros", []),
        })

    # Campaña actual
    actual = _extract_current_campaign(datos, depto)

    # Totales históricos
    total_prima = sum(d["prima_neta"] for d in campanas_data)
    total_indem = sum(d["monto"] for d in campanas_data)
    total_avisos = sum(d["avisos"] for d in campanas_data)
    total_indemnizados = sum(d["indemnizados"] for d in campanas_data)
    total_ha = sum(d["ha_ind"] for d in campanas_data)
    sin_global = round(100 * total_indem / total_prima, 1) if total_prima > 0 else 0

    # ══════════════════════════════════════════════════════════
    # CREAR PRESENTACIÓN
    # ══════════════════════════════════════════════════════════
    prs = Presentation()
    prs.slide_width = Inches(13.33); prs.slide_height = Inches(7.5)

    # SLIDE 1: PORTADA
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, prs, GRIS); _lines(slide, prs)
    _text(slide, 0.8, 1.8, 11.7, 1.0, "SEGURO AGRÍCOLA CATASTRÓFICO", 40, True, BLANCO, PP_ALIGN.CENTER, "Georgia")
    _text(slide, 0.8, 2.7, 11.7, 0.5, "SAC — Análisis Histórico de Siniestralidad", 22, False, TEAL, PP_ALIGN.CENTER, "Georgia")
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.4), Inches(4.3), Inches(0.04))
    sep.fill.solid(); sep.fill.fore_color.rgb = VERDE; sep.line.fill.background()
    _text(slide, 0.8, 3.7, 11.7, 0.5, f"Departamento: {depto_display}", 28, True, BLANCO, PP_ALIGN.CENTER, "Georgia")
    _text(slide, 0.8, 4.4, 11.7, 0.4, "5 campañas agrícolas (2020-2021 a 2024-2025)", 14, False, TEAL, PP_ALIGN.CENTER)
    _text(slide, 0.8, 5.0, 11.7, 0.4,
          f"Siniestralidad acumulada: {sin_global}% "
          f"({_fmt(total_indem)} indemnizados / {_fmt(total_prima)} prima neta)",
          13, False, DORADO, PP_ALIGN.CENTER)
    _text(slide, 0.8, 5.8, 11.7, 0.8,
          "Dirección de Seguro y Fomento del Financiamiento Agrario — MIDAGRI",
          11, True, RGBColor(0xAA, 0xAA, 0xAA), PP_ALIGN.CENTER)

    # SLIDE 2: RESUMEN EJECUTIVO
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, prs, CREMA)
    _text(slide, 0.8, 0.35, 11.7, 0.5, f"Resumen Ejecutivo — {depto_display} (5 campañas)", 26, True, GRIS, PP_ALIGN.LEFT, "Georgia")

    _kpi_card(slide, 0.6, 1.1, 2.4, 1.3, "PRIMA NETA ACUMULADA", _fmt(total_prima), "5 campañas", AZUL)
    _kpi_card(slide, 3.2, 1.1, 2.4, 1.3, "INDEMNIZACIÓN TOTAL", _fmt(total_indem), "5 campañas", DORADO)
    _kpi_card(slide, 5.8, 1.1, 2.4, 1.3, "SINIESTRALIDAD", f"{sin_global}%", "Indemn. / Prima Neta", _sin_color(sin_global))
    _kpi_card(slide, 8.4, 1.1, 2.4, 1.3, "AVISOS TOTALES", f"{total_avisos:,}", f"{total_indemnizados:,} indemnizados", VERDE)
    _kpi_card(slide, 11.0, 1.1, 1.7, 1.3, "HA INDEMNIZADAS", f"{total_ha:,.0f}", "hectáreas", TEAL)

    # Tabla resumen
    _text(slide, 0.6, 2.7, 12.1, 0.35, "Evolución de Siniestralidad por Campaña", 14, True, GRIS, PP_ALIGN.LEFT, "Georgia")
    hdrs = ["Campaña", "Prima Neta", "Indemnización", "Siniestralidad", "Avisos", "Ha Indemn."]
    ws = [1.8, 1.8, 1.8, 1.6, 1.2, 1.2]
    _table_header(slide, 0.6, 3.05, ws, hdrs, GRIS)
    for i, d in enumerate(campanas_data):
        vals = [d["campana"], _fmt(d["prima_neta"]), _fmt(d["monto"]),
                f"{d['siniestralidad']}%", f"{d['avisos']:,}", f"{d['ha_ind']:,.0f}"]
        _table_row(slide, 0.6, 3.4 + i * 0.28, ws, vals, i)

    # SLIDE 3: GRÁFICO PRIMA vs INDEMNIZACIÓN
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, prs, CREMA)
    _text(slide, 0.8, 0.35, 11.7, 0.5, "Prima Neta vs. Indemnización — Siniestralidad", 26, True, GRIS, PP_ALIGN.LEFT, "Georgia")

    chart_data = CategoryChartData()
    chart_data.categories = [d["campana"] for d in campanas_data]
    chart_data.add_series("Prima Neta (S/)", [_safe_chart_val(d["prima_neta"]) / 1e6 for d in campanas_data])
    chart_data.add_series("Indemnización (S/)", [_safe_chart_val(d["monto"]) / 1e6 for d in campanas_data])

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.8), Inches(1.1), Inches(8.0), Inches(4.5), chart_data)
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = AZUL
    chart.series[1].format.fill.solid(); chart.series[1].format.fill.fore_color.rgb = DORADO

    # Panel de siniestralidad
    _text(slide, 9.2, 1.1, 3.5, 0.4, "Siniestralidad (%)", 18, True, GRIS, PP_ALIGN.CENTER, "Georgia")
    for i, d in enumerate(campanas_data):
        y = 1.7 + i * 1.0
        sin_val = _safe_chart_val(d["siniestralidad"])
        bar_w = min(3.2, sin_val / 100 * 3.2)
        bc = _sin_color(sin_val)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      Inches(9.4), Inches(y + 0.35), Inches(max(0.1, bar_w)), Inches(0.3))
        bar.fill.solid(); bar.fill.fore_color.rgb = bc; bar.line.fill.background()
        _text(slide, 9.4, y, 3.2, 0.35, d["campana"], 11, True, GRIS)
        _text(slide, 9.4 + bar_w + 0.15, y + 0.35, 1.0, 0.3, f"{d['siniestralidad']}%", 14, True, bc)

    _obs_box(slide, 0.8, 5.9, 12.0, 0.6,
             "Siniestralidad = Indemnización / Prima Neta × 100%. "
             "Valores > 70% indican campañas con pérdidas que superan ampliamente lo esperado por el seguro.")

    # SLIDES 4-8: DETALLE POR CAMPAÑA HISTÓRICA
    for idx, d in enumerate(campanas_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(slide, prs, CREMA)

        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.3), Inches(2.5), Inches(0.45))
        badge.fill.solid(); badge.fill.fore_color.rgb = VERDE; badge.line.fill.background()
        _text(slide, 0.55, 0.3, 2.4, 0.45, f"  Campaña {d['campana']}", 14, True, BLANCO)
        _text(slide, 3.3, 0.32, 9.5, 0.4, f"{depto_display} — Detalle de Campaña", 20, True, GRIS, PP_ALIGN.LEFT, "Georgia")

        sc = _sin_color(d["siniestralidad"])
        mini = [
            ("Siniestralidad", f"{d['siniestralidad']}%", sc),
            ("Prima Neta", _fmt(d["prima_neta"]), AZUL),
            ("Indemnización", _fmt(d["monto"]), DORADO),
            ("Avisos", f"{d['avisos']:,}", TEAL),
            ("Indemnizados", f"{d['indemnizados']:,}", VERDE),
            ("Ha Indemnizadas", f"{d['ha_ind']:,.0f}", GRIS),
        ]
        for j, (lbl, val, acc) in enumerate(mini):
            _kpi_card(slide, 0.5 + j * 2.1, 0.95, 1.9, 1.0, lbl, val, "", acc)

        # Top cultivos
        _text(slide, 0.5, 2.15, 6.0, 0.35, "Top 5 Cultivos Afectados", 15, True, GRIS, PP_ALIGN.LEFT, "Georgia")
        cult_hdrs = ["Cultivo", "Avisos", "%"]
        cult_ws = [3.3, 1.1, 1.3]
        _table_header(slide, 0.5, 2.5, cult_ws, cult_hdrs, VERDE)
        for ri, cult_item in enumerate(d["top_cultivos"][:5]):
            name, cnt = cult_item[0], cult_item[1]
            pct = f"{100*cnt/max(d['avisos'],1):.1f}%"
            _table_row(slide, 0.5, 2.85 + ri * 0.38, cult_ws, [name.title(), str(cnt), pct], ri)

        # Top siniestros
        _text(slide, 6.8, 2.15, 6.0, 0.35, "Top 5 Tipos de Siniestro", 15, True, GRIS, PP_ALIGN.LEFT, "Georgia")
        sin_hdrs = ["Evento", "Avisos", "%"]
        sin_ws = [2.8, 1.0, 1.0]
        _table_header(slide, 6.8, 2.5, sin_ws, sin_hdrs, TEAL)
        for ri, sin_item in enumerate(d["top_siniestros"][:5]):
            name, cnt = sin_item[0], sin_item[1]
            pct = f"{100*cnt/max(d['avisos'],1):.1f}%"
            _table_row(slide, 6.8, 2.85 + ri * 0.38, sin_ws, [name.title(), str(cnt), pct], ri)

        # Contexto
        diff = d["prima_neta"] - d["monto"]
        estado = "superávit" if diff > 0 else "déficit"
        _text(slide, 0.5, 5.0, 12.3, 0.5,
              f"Prima Neta: {_fmt(d['prima_neta'])} · Indemnización: {_fmt(d['monto'])} · "
              f"Diferencia: {_fmt(abs(diff))} ({estado}) · "
              f"{d['provincias']} provincias · {d['distritos']} distritos",
              11, False, RGBColor(0x64, 0x74, 0x8B), PP_ALIGN.CENTER)

    # SLIDE 9: CAMPAÑA ACTUAL (datos dinámicos)
    if actual:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(slide, prs, CREMA)

        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.3), Inches(3.2), Inches(0.45))
        badge.fill.solid(); badge.fill.fore_color.rgb = DORADO; badge.line.fill.background()
        _text(slide, 0.55, 0.3, 3.1, 0.45, "  Campaña 2025-2026 (en curso)", 14, True, GRIS)
        _text(slide, 4.0, 0.32, 9.0, 0.4, f"{depto_display} — Avance al Corte {actual['fecha_corte']}", 20, True, GRIS, PP_ALIGN.LEFT, "Georgia")

        sc = _sin_color(actual["siniestralidad"])
        mini = [
            ("Siniestralidad", f"{actual['siniestralidad']}%", sc),
            ("Prima Neta", _fmt(actual["prima_neta"]), AZUL),
            ("Indemnización", _fmt(actual["monto_indemnizado"]), DORADO),
            ("Avisos", f"{actual['avisos']:,}", TEAL),
            ("Indemnizados", f"{actual['indemnizados']:,}", VERDE),
            ("Pendientes", f"{actual['pendientes']:,}", NARANJA),
        ]
        for j, (lbl, val, acc) in enumerate(mini):
            _kpi_card(slide, 0.5 + j * 2.1, 0.95, 1.9, 1.0, lbl, val, "", acc)

        # Top cultivos actual
        _text(slide, 0.5, 2.15, 6.0, 0.35, "Top 5 Cultivos Afectados", 15, True, GRIS, PP_ALIGN.LEFT, "Georgia")
        _table_header(slide, 0.5, 2.5, cult_ws, cult_hdrs, VERDE)
        for ri, (name, cnt) in enumerate(actual["top_cultivos"][:5]):
            pct = f"{100*cnt/max(actual['avisos'],1):.1f}%"
            _table_row(slide, 0.5, 2.85 + ri * 0.38, cult_ws, [name.title(), str(cnt), pct], ri)

        # Top siniestros actual
        _text(slide, 6.8, 2.15, 6.0, 0.35, "Top 5 Tipos de Siniestro", 15, True, GRIS, PP_ALIGN.LEFT, "Georgia")
        _table_header(slide, 6.8, 2.5, sin_ws, sin_hdrs, TEAL)
        for ri, (name, cnt) in enumerate(actual["top_siniestros"][:5]):
            pct = f"{100*cnt/max(actual['avisos'],1):.1f}%"
            _table_row(slide, 6.8, 2.85 + ri * 0.38, sin_ws, [name.title(), str(cnt), pct], ri)

        _obs_box(slide, 0.5, 5.0, 12.3, 0.7,
                 f"Campaña en curso con {actual['siniestralidad']}% de siniestralidad al corte "
                 f"({_fmt(actual['monto_indemnizado'])} de {_fmt(actual['prima_neta'])} de prima). "
                 f"Hay {actual['pendientes']} avisos pendientes de dictamen y "
                 f"{actual['no_indemnizables']} no indemnizables. "
                 f"Desembolsado: {_fmt(actual['monto_desembolsado'])}.")

        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.65))
        box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD); box.line.fill.background()
        _text(slide, 0.7, 5.9, 11.9, 0.55,
              f"Referencia histórica: Siniestralidad promedio 5 campañas = {sin_global}% · "
              f"Prima Neta actual {_fmt(actual['prima_neta'])} vs prom. histórico {_fmt(total_prima/5)} · "
              f"{actual['provincias']} provincias, {actual['distritos']} distritos · "
              f"Fuente: Consolidado SAC descargado de aseguradoras al {actual['fecha_corte']}",
              10, False, RGBColor(0x37, 0x47, 0x4F), PP_ALIGN.CENTER)

    # SLIDE 10: TOP PROVINCIAS
    if top_provs:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(slide, prs, CREMA)
        _text(slide, 0.8, 0.35, 11.7, 0.5,
              f"Provincias de Mayor Impacto — {depto_display} (acumulado 5 campañas)", 22, True, GRIS, PP_ALIGN.LEFT, "Georgia")

        hdrs = ["Provincia", "Avisos", "Indemnizados", "Monto Indemnizado", "% del total"]
        ws_p = [2.2, 1.3, 1.5, 2.3, 1.5]
        _table_header(slide, 0.8, 1.0, ws_p, hdrs, GRIS)
        for ri, prov_data in enumerate(top_provs[:10]):
            prov, av, ind, monto = prov_data
            pct = f"{100*monto/max(total_indem,1):.1f}%"
            _table_row(slide, 0.8, 1.35 + ri * 0.42, ws_p, [prov.title(), f"{av:,}", f"{ind:,}", _fmt(monto), pct], ri)

    # SLIDE 11: ESTACIONALIDAD
    if estacionalidad:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(slide, prs, CREMA)
        _text(slide, 0.8, 0.35, 11.7, 0.5,
              f"Estacionalidad de Siniestros — {depto_display}", 22, True, GRIS, PP_ALIGN.LEFT, "Georgia")

        chart_data = CategoryChartData()
        meses = ["Ene", "Feb", "Mar", "Abr", "May", "Jun", "Jul", "Ago", "Sep", "Oct", "Nov", "Dic"]
        chart_data.categories = meses
        chart_data.add_series("Avisos", [estacionalidad.get(str(m), 0) for m in range(1, 13)])

        chart_frame = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.8), Inches(1.0), Inches(11.5), Inches(4.0), chart_data)
        chart = chart_frame.chart
        chart.has_legend = False
        chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = VERDE

        # Identificar mes pico
        if estacionalidad:
            mes_pico = max(estacionalidad.items(), key=lambda x: x[1])
            total_est = sum(estacionalidad.values())
            _text(slide, 0.8, 5.2, 11.5, 0.35, "Interpretación", 16, True, GRIS, PP_ALIGN.LEFT, "Georgia")
            _text(slide, 0.8, 5.55, 11.5, 1.0,
                  f"El mes con mayor actividad es {meses[int(mes_pico[0])-1]} con {mes_pico[1]:,} avisos "
                  f"acumulados en 5 campañas. Los meses de enero a marzo suelen concentrar la mayor "
                  f"actividad por la temporada de lluvias y el desarrollo vegetativo de los cultivos transitorios.",
                  12, False, GRIS)

    # SLIDE 12: CIERRE
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, prs, GRIS); _lines(slide, prs)
    _text(slide, 0.8, 2.0, 11.7, 0.8, "SEGURO AGRÍCOLA CATASTRÓFICO", 32, True, BLANCO, PP_ALIGN.CENTER, "Georgia")
    _text(slide, 0.8, 2.8, 11.7, 0.5, f"SAC 2020–2025 · Departamento: {depto_display}", 18, False, TEAL, PP_ALIGN.CENTER, "Georgia")
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.5), Inches(3.5), Inches(4.3), Inches(0.04))
    sep.fill.solid(); sep.fill.fore_color.rgb = VERDE; sep.line.fill.background()
    _text(slide, 0.8, 3.8, 11.7, 0.8,
          "Dirección de Seguro y Fomento del Financiamiento Agrario\n"
          "Ministerio de Desarrollo Agrario y Riego — MIDAGRI", 13, False, RGBColor(0xAA, 0xAA, 0xAA), PP_ALIGN.CENTER)
    _text(slide, 0.8, 5.0, 11.7, 0.7,
          f"Siniestralidad = Indemnización / Prima Neta = {sin_global}%\n"
          f"Fuente: Datos históricos de 5 campañas SAC · Primas: FOGASA/MIDAGRI\n"
          f"Campaña actual: Consolidado SAC descargado de aseguradoras",
          10, True, RGBColor(0x88, 0x88, 0x88), PP_ALIGN.CENTER)

    # Guardar a bytes
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
